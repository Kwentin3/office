from __future__ import annotations

import copy
import json
import subprocess
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path
from unittest import mock

from lxml import etree
from pptx import Presentation
from pptx.util import Inches
from pptx_artifact_tool import PptxArtifactTool


class StrictTemplateTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temp = tempfile.TemporaryDirectory()
        self.root = Path(self.temp.name)
        self.tool = PptxArtifactTool(self.root / "work")
        self.source = self.root / "template.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        title.name = "slot:title"
        title.text_frame.clear()
        title.text_frame.paragraphs[0].add_run().text = "Offer for {{cus"
        title.text_frame.paragraphs[0].add_run().text = "tomer}}"
        table = slide.shapes.add_table(1, 1, Inches(1), Inches(2), Inches(5), Inches(1))
        table.name = "slot:price"
        table.table.cell(0, 0).text = "{{amount}} {{currency}}"
        presentation.save(self.source)

    def tearDown(self) -> None:
        self.temp.cleanup()

    def test_strict_fill_resolves_managed_slots_and_table_cells(self) -> None:
        output = self.root / "filled.pptx"
        result = self.tool.fill_template(
            self.source,
            {"customer": "Acme", "amount": "100", "currency": "USD"},
            output,
        )
        self.assertEqual(result["status"], "ok")
        summary = self.tool.inspect(output)
        slide = self.tool.inspect(output, view="slide", slide_id=summary["slides"][0]["id"])
        self.assertEqual(slide["slots"][0]["text"], "Offer for Acme")
        self.assertEqual(slide["table_cells"][0]["text"], "100 USD")

    def test_unmanaged_tokens_are_refused_in_strict_mode(self) -> None:
        presentation = Presentation(self.source)
        shape = presentation.slides[0].shapes.add_textbox(Inches(1), Inches(4), Inches(5), Inches(1))
        shape.name = "decorative"
        shape.text = "Hidden {{unsupported}}"
        presentation.save(self.source)
        output = self.root / "never.pptx"
        result = self.tool.fill_template(
            self.source,
            {"customer": "Acme", "amount": "100", "currency": "USD", "unsupported": "x"},
            output,
        )
        self.assertEqual(result["status"], "refused")
        self.assertEqual(result["reason"], "unsupported_capability")
        self.assertFalse(output.exists())

    def test_marker_free_speaker_notes_are_preserved(self) -> None:
        presentation = Presentation(self.source)
        presentation.slides[0].notes_slide.notes_text_frame.text = "Keep note"
        presentation.save(self.source)
        output = self.root / "notes-output.pptx"
        result = self.tool.fill_template(
            self.source,
            {"customer": "Acme", "amount": "100", "currency": "USD"},
            output,
        )
        self.assertEqual(result["status"], "ok")
        self.assertIn("Keep note", Presentation(output).slides[0].notes_slide.notes_text_frame.text)

    def test_strict_fill_rejects_package_wide_marker(self) -> None:
        rewritten = self.root / "package-marker.pptx"
        with zipfile.ZipFile(self.source) as package, zipfile.ZipFile(rewritten, "w") as target:
            for info in package.infolist():
                payload = package.read(info.filename)
                if info.filename == "docProps/core.xml":
                    root = etree.fromstring(payload)
                    root.set("{urn:kwentin:test}marker", "{{residual}}")
                    payload = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                target.writestr(copy.copy(info), payload)
        output = self.root / "package-marker-output.pptx"
        result = self.tool.fill_template(
            rewritten,
            {"customer": "Acme", "amount": "100", "currency": "USD"},
            output,
        )
        self.assertEqual((result["status"], result["reason"]), ("refused", "unsupported_capability"))
        self.assertFalse(output.exists())

    def test_strict_fill_preserves_non_token_runs(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        shape.name = "slot:title"
        paragraph = shape.text_frame.paragraphs[0]
        prefix = paragraph.add_run()
        prefix.text = "Prefix "
        prefix.font.italic = True
        token = paragraph.add_run()
        token.text = "{{name}}"
        token.font.bold = True
        source = self.root / "formatted.pptx"
        presentation.save(source)
        output = self.root / "formatted-output.pptx"
        result = self.tool.fill_template(source, {"name": "Acme"}, output)
        self.assertEqual(result["status"], "ok")
        rendered = Presentation(output).slides[0].shapes[0].text_frame.paragraphs[0]
        self.assertEqual(rendered.text, "Prefix Acme")
        self.assertTrue(rendered.runs[0].font.italic)
        self.assertTrue(rendered.runs[1].font.bold)

    def test_strict_fill_refuses_malformed_overflow_and_group_scope(self) -> None:
        malformed = Presentation()
        slide = malformed.slides.add_slide(malformed.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        shape.name = "slot:title"
        shape.text = "{{{name}}}"
        malformed_source = self.root / "malformed.pptx"
        malformed.save(malformed_source)
        result = self.tool.fill_template(malformed_source, {"name": "Acme"}, self.root / "malformed-output.pptx")
        self.assertEqual((result["status"], result["reason"]), ("refused", "validation_failure"))

        overflow = Presentation()
        slide = overflow.slides.add_slide(overflow.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        shape.name = "slot:title"
        shape.text = "X{{name}}"
        overflow_source = self.root / "overflow.pptx"
        overflow.save(overflow_source)
        result = self.tool.fill_template(overflow_source, {"name": "x" * 32767}, self.root / "overflow-output.pptx")
        self.assertEqual((result["status"], result["reason"]), ("refused", "validation_failure"))

        grouped = Presentation()
        slide = grouped.slides.add_slide(grouped.slide_layouts[6])
        top = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        top.name = "slot:title"
        top.text = "{{name}}"
        group = slide.shapes.add_group_shape()
        hidden = group.shapes.add_textbox(Inches(1), Inches(3), Inches(5), Inches(1))
        hidden.name = "slot:hidden"
        hidden.text_frame.clear()
        hidden.text_frame.paragraphs[0].add_run().text = "{{hid"
        hidden.text_frame.paragraphs[0].add_run().text = "den}}"
        grouped_source = self.root / "grouped.pptx"
        grouped.save(grouped_source)
        result = self.tool.fill_template(
            grouped_source, {"name": "Acme", "hidden": "x"}, self.root / "grouped-output.pptx"
        )
        self.assertEqual((result["status"], result["reason"]), ("refused", "unsupported_capability"))

    def test_strict_fill_rejects_residual_marker_in_unsupported_candidate_scope(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        shape.name = "slot:title"
        shape.text = "{{name}}"
        source = self.root / "residual-source.pptx"
        presentation.save(source)
        output = self.root / "residual-output.pptx"
        original_apply = self.tool.apply

        def inject_residual(source_path, plan, candidate):
            result = original_apply(source_path, plan, candidate)
            if result.get("status") == "ok":
                rewritten = self.root / "residual-candidate.pptx"
                with zipfile.ZipFile(candidate) as package, zipfile.ZipFile(rewritten, "w") as target:
                    for info in package.infolist():
                        payload = package.read(info.filename)
                        if info.filename == "docProps/core.xml":
                            root = etree.fromstring(payload)
                            root.set("{urn:kwentin:test}residual", "{{residual}}")
                            payload = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                        target.writestr(copy.copy(info), payload)
                rewritten.replace(candidate)
            return result

        with mock.patch.object(self.tool, "apply", side_effect=inject_residual):
            result = self.tool.fill_template(source, {"name": "Acme"}, output)
        self.assertEqual((result["status"], result["reason"]), ("refused", "validation_failure"))
        self.assertFalse(output.exists())

    def test_strict_fill_rejects_dtd_entity_declarations_package_wide(self) -> None:
        declaration = '<!DOCTYPE cp:coreProperties [<!ENTITY hidden "{{hidden}}">]>'
        for encoding in ("UTF-8", "UTF-16"):
            with self.subTest(encoding=encoding):
                rewritten = self.root / f"dtd-marker-{encoding}.pptx"
                with zipfile.ZipFile(self.source) as package, zipfile.ZipFile(rewritten, "w") as target:
                    for info in package.infolist():
                        payload = package.read(info.filename)
                        if info.filename == "docProps/core.xml":
                            root = etree.fromstring(payload)
                            payload = etree.tostring(
                                root,
                                xml_declaration=True,
                                encoding=encoding,
                                doctype=declaration,
                            )
                        target.writestr(copy.copy(info), payload)
                output = self.root / f"dtd-output-{encoding}.pptx"
                result = self.tool.fill_template(
                    rewritten,
                    {"customer": "Acme", "amount": "100", "currency": "USD"},
                    output,
                )
                self.assertEqual(result["status"], "refused")
                self.assertFalse(output.exists())

    def test_strict_fill_preserves_managed_multi_paragraph_runs(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
        shape.name = "slot:body"
        first = shape.text_frame.paragraphs[0]
        first_run = first.add_run()
        first_run.text = "First paragraph"
        first_run.font.italic = True
        second = shape.text_frame.add_paragraph()
        prefix = second.add_run()
        prefix.text = "Second prefix "
        prefix.font.italic = True
        token = second.add_run()
        token.text = "{{name}}"
        token.font.bold = True
        source = self.root / "multi-paragraph.pptx"
        presentation.save(source)
        output = self.root / "multi-paragraph-output.pptx"
        result = self.tool.fill_template(source, {"name": "Acme"}, output)
        self.assertEqual(result["status"], "ok")
        rendered = Presentation(output).slides[0].shapes[0].text_frame
        self.assertEqual([p.text for p in rendered.paragraphs], ["First paragraph", "Second prefix Acme"])
        self.assertTrue(rendered.paragraphs[0].runs[0].font.italic)
        self.assertTrue(rendered.paragraphs[1].runs[0].font.italic)
        self.assertTrue(rendered.paragraphs[1].runs[1].font.bold)

    def test_cli_fill_template_matches_api_contract(self) -> None:
        output = self.root / "cli.pptx"
        payload = {
            "action": "fill_template",
            "workdir": str(self.root / "cli-work"),
            "source": str(self.source),
            "values": {"customer": "Acme", "amount": "100", "currency": "USD"},
            "output": str(output),
        }
        completed = subprocess.run(
            [sys.executable, "-m", "pptx_artifact_tool"],
            input=json.dumps(payload),
            text=True,
            capture_output=True,
            check=False,
        )
        self.assertEqual(completed.returncode, 0, completed.stderr)
        self.assertEqual(json.loads(completed.stdout)["status"], "ok")
        self.assertTrue(output.is_file())


if __name__ == "__main__":
    unittest.main()
