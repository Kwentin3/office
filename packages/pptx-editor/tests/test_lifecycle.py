from __future__ import annotations
import tempfile,unittest,zipfile,hashlib
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pptx_artifact_tool import PptxArtifactTool

class LifecycleTests(unittest.TestCase):
 def setUp(self):
  self.temp=tempfile.TemporaryDirectory();self.root=Path(self.temp.name);self.tool=PptxArtifactTool(self.root/'work');self.template=self.root/'template.pptx'
  prs=Presentation();slide=prs.slides.add_slide(prs.slide_layouts[5]);slide.shapes.title.text='{{title}}';slide.shapes.title.name='slot:cover.title';body=slide.shapes.add_textbox(Inches(1),Inches(2),Inches(8),Inches(1));body.name='slot:cover.body';body.text='{{body}}'
  slide2=prs.slides.add_slide(prs.slide_layouts[5]);slide2.shapes.title.text='Results';slide2.shapes.title.name='slot:results.title';table=slide2.shapes.add_table(2,2,Inches(1),Inches(2),Inches(8),Inches(2));table.name='slot:results.table';table.table.cell(0,0).text='Metric';table.table.cell(0,1).text='Value';table.table.cell(1,0).text='Revenue';table.table.cell(1,1).text='{{value}}';prs.save(self.template)
 def tearDown(self):self.temp.cleanup()
 def test_inspect_plan_apply_validate_text_and_table_cells(self):
  summary=self.tool.inspect(self.template,view='summary');self.assertEqual(summary['status'],'ok');self.assertEqual(len(summary['slides']),2);self.assertEqual({x['key'] for s in summary['slides'] for x in s['slots']},{'cover.title','cover.body','results.title','results.table'})
  slide=self.tool.inspect(self.template,view='slide',slide_id=summary['slides'][0]['id']);title=next(x for x in slide['slots'] if x['key']=='cover.title');body=next(x for x in slide['slots'] if x['key']=='cover.body')
  slide2=self.tool.inspect(self.template,view='slide',slide_id=summary['slides'][1]['id']);cell=next(x for x in slide2['table_cells'] if x['slot_key']=='results.table' and x['row']==1 and x['column']==1)
  request={'operations':[{'type':'set_slot_text','target_id':title['id'],'text':'Quarterly Review','expected_text':'{{title}}'},{'type':'set_slot_text','target_id':body['id'],'text':'Revenue up\nCosts stable','expected_text':'{{body}}'},{'type':'set_table_cell_text','target_id':cell['id'],'text':'125','expected_text':'{{value}}'}]}
  planned=self.tool.plan({'summary':summary,'slides':[slide,slide2]},request);self.assertEqual(planned['status'],'ok');output=self.root/'output.pptx';applied=self.tool.apply(self.template,planned['plan'],output);self.assertEqual(applied['status'],'ok')
  after=Presentation(output);self.assertEqual(after.slides[0].shapes.title.text,'Quarterly Review');self.assertEqual(next(x for x in after.slides[0].shapes if x.name=='slot:cover.body').text,'Revenue up\vCosts stable');self.assertEqual(next(x for x in after.slides[1].shapes if x.has_table).table.cell(1,1).text,'125')
  report=self.tool.validate(output,before=self.template);self.assertEqual(report['status'],'valid');self.assertEqual(report['unexpected_changed_members'],[])
 def test_create_renders_stable_slot_keys(self):
  output=self.root/'rendered.pptx';result=self.tool.create(self.template,{'slots':{'cover.title':'Cake Plan','cover.body':'Ingredients\nBudget','results.title':'Totals'},'table_cells':{'results.table.r1c1':'500'}},output);self.assertEqual(result['status'],'ok');after=Presentation(output);self.assertEqual(after.slides[0].shapes.title.text,'Cake Plan');self.assertEqual(next(x for x in after.slides[1].shapes if x.has_table).table.cell(1,1).text,'500')

 def test_duplicate_template_slot_keys_refuse_create(self):
  prs=Presentation(self.template);extra=prs.slides[1].shapes.add_textbox(Inches(1),Inches(5),Inches(4),Inches(1));extra.name='slot:cover.title';extra.text='duplicate';duplicate=self.root/'duplicate.pptx';prs.save(duplicate);output=self.root/'must-not-exist.pptx'
  result=self.tool.create(duplicate,{'slots':{'cover.title':'New'}},output);self.assertEqual(result['status'],'refused');self.assertEqual(result['reason'],'ambiguous_target');self.assertFalse(output.exists())

 def test_edit_after_reorder_targets_logical_slide_relationship(self):
  summary=self.tool.inspect(self.template,view='summary');slides=[self.tool.inspect(self.template,view='slide',slide_id=s['id']) for s in summary['slides']];reorder=self.tool.plan({'summary':summary,'slides':slides},{'operations':[{'type':'reorder_slides','slide_ids':[summary['slides'][1]['id'],summary['slides'][0]['id']]}]})['plan'];reordered=self.root/'reordered-then-edit.pptx';self.assertEqual(self.tool.apply(self.template,reorder,reordered)['status'],'ok')
  summary2=self.tool.inspect(reordered,view='summary');first=self.tool.inspect(reordered,view='slide',slide_id=summary2['slides'][0]['id']);title=next(x for x in first['slots'] if x['key']=='results.title');plan=self.tool.plan({'summary':summary2,'slides':[first]},{'operations':[{'type':'set_slot_text','target_id':title['id'],'text':'Changed after reorder','expected_text':'Results'}]})['plan'];output=self.root/'after-reorder-edit.pptx';self.assertEqual(self.tool.apply(reordered,plan,output)['status'],'ok');after=Presentation(output);self.assertEqual(after.slides[0].shapes.title.text,'Changed after reorder');self.assertEqual(after.slides[1].shapes.title.text,'{{title}}')

 def test_slot_after_picture_is_addressed_by_powerpoint_shape_id(self):
  from PIL import Image
  image=self.root/'pixel.png';Image.new('RGB',(8,8),(255,0,0)).save(image)
  prs=Presentation();slide=prs.slides.add_slide(prs.slide_layouts[6]);slide.shapes.add_picture(str(image),Inches(0.2),Inches(0.2),Inches(0.5),Inches(0.5));box=slide.shapes.add_textbox(Inches(1),Inches(1),Inches(5),Inches(1));box.name='slot:body';box.text='Old';template=self.root/'picture-template.pptx';prs.save(template)
  output=self.root/'picture-output.pptx';result=self.tool.create(template,{'slots':{'body':'New'}},output);self.assertEqual(result['status'],'ok');after=Presentation(output);self.assertEqual(next(x for x in after.slides[0].shapes if x.name=='slot:body').text,'New');self.assertEqual(len([x for x in after.slides[0].shapes if x.shape_type.name=='PICTURE']),1)

 def test_multiline_slot_uses_powerpoint_breaks(self):
  output=self.root/'multiline.pptx';result=self.tool.create(self.template,{'slots':{'cover.body':'Line one\nLine two'}},output);self.assertEqual(result['status'],'ok')
  with zipfile.ZipFile(output) as archive:xml=archive.read('ppt/slides/slide1.xml')
  self.assertIn(b'<a:br',xml);after=Presentation(output);self.assertEqual(next(x for x in after.slides[0].shapes if x.name=='slot:cover.body').text,'Line one\vLine two')

 def test_reorder_slides_changes_only_presentation_order(self):
  summary=self.tool.inspect(self.template,view='summary');slides=[self.tool.inspect(self.template,view='slide',slide_id=s['id']) for s in summary['slides']]
  planned=self.tool.plan({'summary':summary,'slides':slides},{'operations':[{'type':'reorder_slides','slide_ids':[summary['slides'][1]['id'],summary['slides'][0]['id']]}]});self.assertEqual(planned['status'],'ok')
  output=self.root/'reordered.pptx';result=self.tool.apply(self.template,planned['plan'],output);self.assertEqual(result['status'],'ok');after=Presentation(output);self.assertEqual([s.shapes.title.text for s in after.slides],['Results','{{title}}'])
  report=self.tool.validate(output,before=self.template);self.assertEqual(report['status'],'valid');self.assertEqual(report['changed_members'],['ppt/presentation.xml'])

if __name__=='__main__':unittest.main()
