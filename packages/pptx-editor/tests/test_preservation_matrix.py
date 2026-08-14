from __future__ import annotations

import tempfile
import unittest
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from pptx_artifact_tool import PptxArtifactTool


class PptxPreservationMatrixTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temp = tempfile.TemporaryDirectory()
        self.root = Path(self.temp.name)
        self.tool = PptxArtifactTool(self.root / "work")

    def tearDown(self) -> None:
        self.temp.cleanup()

    @staticmethod
    def members(path: Path) -> dict[str, bytes]:
        with zipfile.ZipFile(path) as archive:
            return {info.filename: archive.read(info.filename) for info in archive.infolist()}

    def deck(self) -> Path:
        source = self.root / "rich.pptx"
        presentation = Presentation()
        first = presentation.slides.add_slide(presentation.slide_layouts[6])
        linked = first.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        linked.name = "slot:first"
        run = linked.text_frame.paragraphs[0].add_run()
        run.text = "Linked label"
        run.hyperlink.address = "https://example.invalid/preserve"
        first.notes_slide.notes_text_frame.text = "Preserve speaker note"
        second = presentation.slides.add_slide(presentation.slide_layouts[6])
        target = second.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        target.name = "slot:second"
        target.text = "Old second"
        table = second.shapes.add_table(1, 1, Inches(1), Inches(3), Inches(5), Inches(1))
        table.name = "slot:table"
        table.table.cell(0, 0).text = "Old cell"
        presentation.save(source)
        return source

    def plan_slot(self, source: Path, slide_index: int, key: str, text: str) -> dict:
        summary = self.tool.inspect(source, view="summary")
        slide = self.tool.inspect(source, view="slide", slide_id=summary["slides"][slide_index]["id"])
        target = next(item for item in slide["slots"] if item["key"] == key)
        planned = self.tool.plan(
            {"summary": summary, "slides": [self.tool.inspect(source, view="slide", slide_id=item["id"]) for item in summary["slides"]]},
            {"operations": [{"type": "set_slot_text", "target_id": target["id"], "text": text, "expected_text": target["text"]}]},
        )
        self.assertEqual(planned["status"], "ok")
        return planned["plan"]

    def test_different_slide_edit_byte_preserves_hyperlink_relationship_and_notes(self) -> None:
        source = self.deck()
        before = self.members(source)
        output = self.root / "different-slide.pptx"
        result = self.tool.apply(source, self.plan_slot(source, 1, "second", "Changed second"), output)
        self.assertEqual(result["status"], "ok")
        after = self.members(output)
        self.assertEqual(set(after), set(before))
        changed = {name for name in before if before[name] != after[name]}
        self.assertEqual(changed, {"ppt/slides/slide2.xml"})
        self.assertEqual(after["ppt/slides/slide1.xml"], before["ppt/slides/slide1.xml"])
        self.assertEqual(after["ppt/slides/_rels/slide1.xml.rels"], before["ppt/slides/_rels/slide1.xml.rels"])
        self.assertEqual(after["ppt/notesSlides/notesSlide1.xml"], before["ppt/notesSlides/notesSlide1.xml"])

    def test_same_slide_managed_edit_preserves_hyperlink_and_notes_semantics(self) -> None:
        source = self.deck()
        before = self.members(source)
        output = self.root / "same-slide.pptx"
        result = self.tool.apply(source, self.plan_slot(source, 0, "first", "Changed link label"), output)
        self.assertEqual(result["status"], "ok")
        after = self.members(output)
        changed = {name for name in before if before[name] != after[name]}
        self.assertEqual(changed, {"ppt/slides/slide1.xml"})
        self.assertEqual(after["ppt/slides/_rels/slide1.xml.rels"], before["ppt/slides/_rels/slide1.xml.rels"])
        self.assertEqual(after["ppt/notesSlides/notesSlide1.xml"], before["ppt/notesSlides/notesSlide1.xml"])
        presentation = Presentation(output)
        shape = next(item for item in presentation.slides[0].shapes if item.name == "slot:first")
        self.assertEqual(shape.text, "Changed link label")
        self.assertEqual(shape.text_frame.paragraphs[0].runs[0].hyperlink.address, "https://example.invalid/preserve")
        self.assertIn("Preserve speaker note", presentation.slides[0].notes_slide.notes_text_frame.text)

    def test_slide_reorder_preserves_every_slide_and_relationship_member(self) -> None:
        source = self.deck()
        before = self.members(source)
        summary = self.tool.inspect(source, view="summary")
        plan = self.tool.plan(
            {"summary": summary, "slides": [self.tool.inspect(source, view="slide", slide_id=item["id"]) for item in summary["slides"]]},
            {"operations": [{"type": "reorder_slides", "slide_ids": [summary["slides"][1]["id"], summary["slides"][0]["id"]]}]},
        )["plan"]
        output = self.root / "reordered.pptx"
        result = self.tool.apply(source, plan, output)
        self.assertEqual(result["status"], "ok")
        after = self.members(output)
        self.assertEqual({name for name in before if before[name] != after[name]}, {"ppt/presentation.xml"})
        for name in before:
            if name.startswith(("ppt/slides/", "ppt/notesSlides/")):
                self.assertEqual(after[name], before[name], name)


if __name__ == "__main__":
    unittest.main()
