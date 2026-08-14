from __future__ import annotations

import copy
import json
import subprocess
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Inches
from pptx_artifact_tool import PptxArtifactTool
from pptx_artifact_tool.inventory import BLOCKER_FEATURES, WARNING_FEATURES

PPTX_FEATURES = {
    "slides",
    "managed_slots",
    "duplicate_slot_keys",
    "group_shapes",
    "table_shapes",
    "charts",
    "smartart",
    "media",
    "speaker_notes",
    "animations",
    "transitions",
    "hyperlinks",
    "external_relationships",
    "comments",
    "embedded_packages",
    "ole_objects",
    "activex",
    "modification_protection",
    "macros",
    "signatures",
}


class RichInventoryTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temp = tempfile.TemporaryDirectory()
        self.root = Path(self.temp.name)
        self.tool = PptxArtifactTool(self.root / "work")
        self.source = self.root / "source.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        shape.name = "slot:title"
        shape.text = "Normal"
        presentation.save(self.source)

    def tearDown(self) -> None:
        self.temp.cleanup()

    def add_member(self, name: str, payload: bytes = b"x") -> None:
        rewritten = self.root / "rewritten.pptx"
        with zipfile.ZipFile(self.source) as source, zipfile.ZipFile(rewritten, "w") as output:
            for info in source.infolist():
                output.writestr(copy.copy(info), source.read(info.filename))
            output.writestr(name, payload)
        rewritten.replace(self.source)

    def add_ole_node(self) -> None:
        rewritten = self.root / "rewritten.pptx"
        namespace = "http://schemas.openxmlformats.org/presentationml/2006/main"
        with zipfile.ZipFile(self.source) as source, zipfile.ZipFile(rewritten, "w") as output:
            for info in source.infolist():
                payload = source.read(info.filename)
                if info.filename == "ppt/slides/slide1.xml":
                    root = etree.fromstring(payload)
                    tree = root.find(f".//{{{namespace}}}spTree")
                    etree.SubElement(tree, f"{{{namespace}}}oleObj")
                    payload = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                output.writestr(copy.copy(info), payload)
        rewritten.replace(self.source)

    def test_inventory_reports_safe_baseline(self) -> None:
        result = self.tool.inspect(self.source, view="inventory")
        self.assertEqual(result["status"], "ok")
        self.assertEqual(result["schema_version"], 1)
        self.assertEqual(set(result["features"]), PPTX_FEATURES)
        self.assertEqual(result["mutation_policy"]["decision"], "safe")
        self.assertEqual(result["features"]["ole_objects"], 0)
        self.assertEqual(result["features"]["charts"], 0)

    def test_inventory_schema_matches_policy_partitions(self) -> None:
        schema = json.loads(
            (Path(__file__).parents[1] / "pptx_artifact_tool/resources/inventory.schema.json").read_text()
        )
        policy = schema["properties"]["mutation_policy"]
        self.assertEqual(set(policy["properties"]["blockers"]["items"]["enum"]), set(BLOCKER_FEATURES))
        self.assertEqual(set(policy["properties"]["warnings"]["items"]["enum"]), set(WARNING_FEATURES))
        constraints = {
            rule["if"]["properties"]["decision"]["const"]: rule["then"]["properties"] for rule in policy["allOf"]
        }
        self.assertEqual(constraints["safe"]["blockers"]["maxItems"], 0)
        self.assertEqual(constraints["safe"]["warnings"]["maxItems"], 0)
        self.assertEqual(constraints["safe_with_warnings"]["warnings"]["minItems"], 1)
        self.assertEqual(constraints["refuse_mutation"]["blockers"]["minItems"], 1)

    def test_malformed_package_xml_is_typed_for_inventory_and_apply(self) -> None:
        summary = self.tool.inspect(self.source)
        slide = self.tool.inspect(self.source, view="slide", slide_id=summary["slides"][0]["id"])
        target = slide["slots"][0]
        planned = self.tool.plan(
            {"summary": summary, "slides": [slide]},
            {
                "operations": [
                    {"type": "set_slot_text", "target_id": target["id"], "text": "Changed", "expected_text": "Normal"}
                ]
            },
        )
        rewritten = self.root / "malformed-package.pptx"
        with zipfile.ZipFile(self.source) as source, zipfile.ZipFile(rewritten, "w") as output:
            for info in source.infolist():
                payload = b"<broken" if info.filename == "docProps/core.xml" else source.read(info.filename)
                output.writestr(copy.copy(info), payload)
        rewritten.replace(self.source)
        inventory = self.tool.inspect(self.source, view="inventory")
        self.assertEqual((inventory["status"], inventory["reason"]), ("refused", "validation_failure"))
        result = self.tool.apply(self.source, planned["plan"], self.root / "out.pptx")
        self.assertEqual((result["status"], result["reason"]), ("refused", "validation_failure"))

    def test_ole_objects_are_inventory_blocker_and_apply_refuses(self) -> None:
        summary = self.tool.inspect(self.source)
        slide = self.tool.inspect(self.source, view="slide", slide_id=summary["slides"][0]["id"])
        target = slide["slots"][0]
        planned = self.tool.plan(
            {"summary": summary, "slides": [slide]},
            {
                "operations": [
                    {"type": "set_slot_text", "target_id": target["id"], "text": "Changed", "expected_text": "Normal"}
                ]
            },
        )
        self.add_ole_node()
        inventory = self.tool.inspect(self.source, view="inventory")
        self.assertEqual(inventory["features"]["ole_objects"], 1)
        self.assertEqual(inventory["mutation_policy"]["decision"], "refuse_mutation")
        result = self.tool.apply(self.source, planned["plan"], self.root / "out.pptx")
        self.assertEqual(result["status"], "refused")
        self.assertEqual(result["reason"], "unsupported_capability")

    def test_embedded_package_is_warning_not_ole_blocker(self) -> None:
        self.add_member("ppt/embeddings/chart-workbook.xlsx")
        inventory = self.tool.inspect(self.source, view="inventory")
        self.assertEqual(inventory["features"]["embedded_packages"], 1)
        self.assertEqual(inventory["features"]["ole_objects"], 0)
        self.assertEqual(inventory["mutation_policy"]["decision"], "safe_with_warnings")

    def test_modification_protection_is_blocker(self) -> None:
        rewritten = self.root / "protected.pptx"
        namespace = "http://schemas.openxmlformats.org/presentationml/2006/main"
        with zipfile.ZipFile(self.source) as source, zipfile.ZipFile(rewritten, "w") as output:
            for info in source.infolist():
                payload = source.read(info.filename)
                if info.filename == "ppt/presentation.xml":
                    root = etree.fromstring(payload)
                    etree.SubElement(root, f"{{{namespace}}}modifyVerifier")
                    payload = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                output.writestr(copy.copy(info), payload)
        rewritten.replace(self.source)
        inventory = self.tool.inspect(self.source, view="inventory")
        self.assertEqual(inventory["features"]["modification_protection"], 1)
        self.assertEqual(inventory["mutation_policy"]["decision"], "refuse_mutation")

    def test_malformed_inventory_xml_is_typed_for_api_and_cli(self) -> None:
        rewritten = self.root / "malformed.pptx"
        with zipfile.ZipFile(self.source) as source, zipfile.ZipFile(rewritten, "w") as output:
            for info in source.infolist():
                payload = b"<broken" if info.filename == "ppt/slides/slide1.xml" else source.read(info.filename)
                output.writestr(copy.copy(info), payload)
        rewritten.replace(self.source)
        result = self.tool.inspect(self.source, view="inventory")
        self.assertEqual((result["status"], result["reason"]), ("refused", "validation_failure"))
        completed = subprocess.run(
            [sys.executable, "-m", "pptx_artifact_tool"],
            input=json.dumps({"action": "inspect", "source": str(self.source), "view": "inventory"}),
            text=True,
            capture_output=True,
            check=False,
        )
        self.assertEqual(completed.returncode, 0)
        self.assertEqual(json.loads(completed.stdout)["reason"], "validation_failure")
        self.assertNotIn("Traceback", completed.stderr)


if __name__ == "__main__":
    unittest.main()
