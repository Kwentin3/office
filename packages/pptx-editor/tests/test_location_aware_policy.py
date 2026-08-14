from __future__ import annotations

import copy
import json
import tempfile
import unittest
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Inches
from pptx_artifact_tool import PptxArtifactTool


class LocationAwareInventoryTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temp = tempfile.TemporaryDirectory()
        self.root = Path(self.temp.name)
        self.tool = PptxArtifactTool(self.root / "work")

    def tearDown(self) -> None:
        self.temp.cleanup()

    def rewrite(self, source: Path, changes: dict[str, bytes], additions: dict[str, bytes]) -> None:
        rewritten = self.root / "rewritten.pptx"
        with zipfile.ZipFile(source) as package, zipfile.ZipFile(rewritten, "w") as output:
            for info in package.infolist():
                output.writestr(copy.copy(info), changes.get(info.filename, package.read(info.filename)))
            for name, payload in additions.items():
                output.writestr(name, payload)
        rewritten.replace(source)

    def test_inventory_v2_reports_deterministic_shape_and_relationship_findings(self) -> None:
        source = self.root / "locations.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        shape.name = "slot:title"
        shape.text = "Old"
        slide.notes_slide.notes_text_frame.text = "Keep this marker-free note"
        presentation.save(source)

        first = self.tool.inspect(source, view="inventory")
        second = self.tool.inspect(source, view="inventory")

        self.assertEqual(first, second)
        self.assertEqual(first["schema_version"], 2)
        self.assertFalse(first["findings_truncated"])
        managed = [finding for finding in first["findings"] if finding["feature"] == "managed_slots"]
        self.assertEqual(
            managed,
            [
                {
                    "feature": "managed_slots",
                    "scope": "shape",
                    "part": "ppt/slides/slide1.xml",
                    "slide": {"index": 0, "part": "ppt/slides/slide1.xml"},
                    "shape": {"id": shape.shape_id, "name": "slot:title"},
                }
            ],
        )
        notes = [finding for finding in first["findings"] if finding["feature"] == "speaker_notes"]
        self.assertEqual(
            notes,
            [
                {
                    "feature": "speaker_notes",
                    "scope": "relationship",
                    "part": "ppt/notesSlides/notesSlide1.xml",
                    "slide": {"index": 0, "part": "ppt/slides/slide1.xml"},
                }
            ],
        )

    def test_inventory_classifies_macro_signature_protection_activex_and_ole_as_global(self) -> None:
        source = self.root / "global-features.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        shape.name = "slot:title"
        shape.text = "Old"
        presentation.save(source)

        with zipfile.ZipFile(source) as package:
            presentation_root = etree.fromstring(package.read("ppt/presentation.xml"))
            etree.SubElement(
                presentation_root,
                "{http://schemas.openxmlformats.org/presentationml/2006/main}modifyVerifier",
            )
            slide_root = etree.fromstring(package.read("ppt/slides/slide1.xml"))
            shape_tree = slide_root.find(
                ".//{http://schemas.openxmlformats.org/presentationml/2006/main}spTree"
            )
            etree.SubElement(
                shape_tree,
                "{http://schemas.openxmlformats.org/presentationml/2006/main}oleObj",
            )
            rels_root = etree.fromstring(package.read("ppt/slides/_rels/slide1.xml.rels"))
            etree.SubElement(
                rels_root,
                "{http://schemas.openxmlformats.org/package/2006/relationships}Relationship",
                Id="rIdActiveX",
                Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/activeXControl",
                Target="../activeX/activeX1.xml",
            )
        self.rewrite(
            source,
            {
                "ppt/presentation.xml": etree.tostring(presentation_root, xml_declaration=True, encoding="UTF-8"),
                "ppt/slides/slide1.xml": etree.tostring(slide_root, xml_declaration=True, encoding="UTF-8"),
                "ppt/slides/_rels/slide1.xml.rels": etree.tostring(rels_root, xml_declaration=True, encoding="UTF-8"),
            },
            {
                "ppt/vbaProject.bin": b"macro",
                "_xmlsignatures/sig1.xml": b"<Signature/>",
                "ppt/activeX/activeX1.xml": b"<activeX/>",
            },
        )

        inventory = self.tool.inspect(source, view="inventory")

        self.assertEqual(inventory["mutation_policy"]["decision"], "refuse_mutation")
        for feature in ("macros", "signatures", "modification_protection", "activex", "ole_objects"):
            with self.subTest(feature=feature):
                findings = [item for item in inventory["findings"] if item["feature"] == feature]
                self.assertTrue(findings)
                self.assertEqual({item["scope"] for item in findings}, {"global"})
                self.assertTrue(all(not str(item["part"]).startswith("/") for item in findings))

    def test_relationship_target_cannot_escape_package_root(self) -> None:
        source = self.root / "escape.pptx"
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.save(source)
        with zipfile.ZipFile(source) as package:
            rels_root = etree.fromstring(package.read("ppt/slides/_rels/slide1.xml.rels"))
        etree.SubElement(
            rels_root,
            "{http://schemas.openxmlformats.org/package/2006/relationships}Relationship",
            Id="rIdEscape",
            Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart",
            Target="../../../escape.xml",
        )
        self.rewrite(
            source,
            {"ppt/slides/_rels/slide1.xml.rels": etree.tostring(rels_root, xml_declaration=True, encoding="UTF-8")},
            {},
        )
        result = self.tool.inspect(source, view="inventory")
        self.assertEqual((result["status"], result["reason"]), ("refused", "validation_failure"))

    def test_inventory_locates_slide_shape_and_relationship_preservation_features(self) -> None:
        source = self.root / "preservation-features.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        linked = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        linked.name = "slot:linked"
        run = linked.text_frame.paragraphs[0].add_run()
        run.text = "Linked"
        run.hyperlink.address = "https://example.invalid/keep"
        table = slide.shapes.add_table(1, 1, Inches(1), Inches(3), Inches(4), Inches(1))
        table.name = "slot:data"
        presentation.save(source)

        p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        pr = "http://schemas.openxmlformats.org/package/2006/relationships"
        with zipfile.ZipFile(source) as package:
            slide_root = etree.fromstring(package.read("ppt/slides/slide1.xml"))
            etree.SubElement(slide_root, f"{{{p}}}transition", spd="slow")
            timing = etree.SubElement(slide_root, f"{{{p}}}timing")
            etree.SubElement(timing, f"{{{p}}}tnLst")
            rels_root = etree.fromstring(package.read("ppt/slides/_rels/slide1.xml.rels"))
            for ident, relation_type, target in (
                ("rIdChart", "chart", "../charts/chart1.xml"),
                ("rIdMedia", "image", "../media/image1.png"),
                ("rIdSmartArt", "diagramData", "../diagrams/data1.xml"),
                ("rIdComments", "comments", "../comments/comment1.xml"),
            ):
                etree.SubElement(
                    rels_root,
                    f"{{{pr}}}Relationship",
                    Id=ident,
                    Type=f"http://schemas.openxmlformats.org/officeDocument/2006/relationships/{relation_type}",
                    Target=target,
                )
        self.rewrite(
            source,
            {
                "ppt/slides/slide1.xml": etree.tostring(slide_root, xml_declaration=True, encoding="UTF-8"),
                "ppt/slides/_rels/slide1.xml.rels": etree.tostring(rels_root, xml_declaration=True, encoding="UTF-8"),
            },
            {
                "ppt/charts/chart1.xml": b"<chart/>",
                "ppt/media/image1.png": b"image",
                "ppt/diagrams/data1.xml": b"<data/>",
                "ppt/comments/comment1.xml": b'<p:cmLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cm/></p:cmLst>',
            },
        )

        inventory = self.tool.inspect(source, view="inventory")
        by_feature = {}
        for finding in inventory["findings"]:
            by_feature.setdefault(finding["feature"], []).append(finding)

        self.assertEqual(by_feature["animations"][0]["scope"], "slide")
        self.assertEqual(by_feature["transitions"][0]["scope"], "slide")
        self.assertEqual(by_feature["hyperlinks"][0]["scope"], "shape")
        self.assertEqual(by_feature["hyperlinks"][0]["shape"]["id"], linked.shape_id)
        self.assertEqual(by_feature["table_shapes"][0]["shape"]["id"], table.shape_id)
        for feature in ("charts", "media", "smartart", "comments"):
            with self.subTest(feature=feature):
                self.assertEqual(by_feature[feature][0]["scope"], "relationship")
                self.assertEqual(by_feature[feature][0]["slide"], {"index": 0, "part": "ppt/slides/slide1.xml"})

    def test_inventory_schema_v2_closes_and_bounds_location_findings(self) -> None:
        schema_path = Path(__file__).parents[1] / "pptx_artifact_tool/resources/inventory.schema.json"
        schema = json.loads(schema_path.read_text())

        self.assertEqual(schema["properties"]["schema_version"]["const"], 2)
        self.assertTrue({"findings", "findings_truncated"}.issubset(schema["required"]))
        findings = schema["properties"]["findings"]
        self.assertEqual(findings["maxItems"], 1000)
        item = findings["items"]
        self.assertFalse(item["additionalProperties"])
        self.assertEqual(set(item["required"]), {"feature", "scope", "part"})
        self.assertEqual(
            set(item["properties"]["scope"]["enum"]),
            {"global", "presentation", "slide", "shape", "relationship"},
        )
        self.assertFalse(item["properties"]["slide"]["additionalProperties"])
        self.assertFalse(item["properties"]["shape"]["additionalProperties"])


if __name__ == "__main__":
    unittest.main()
