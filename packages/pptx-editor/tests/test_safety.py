from __future__ import annotations
import copy,tempfile,unittest,zipfile
from pathlib import Path
from unittest.mock import patch
from pptx import Presentation
from pptx.util import Inches
from pptx_artifact_tool import PptxArtifactTool
from pptx_artifact_tool.api import _obj

class SafetyTests(unittest.TestCase):
 def setUp(self):
  self.temp=tempfile.TemporaryDirectory();self.root=Path(self.temp.name);self.tool=PptxArtifactTool(self.root/'work');self.source=self.root/'source.pptx';prs=Presentation();slide=prs.slides.add_slide(prs.slide_layouts[5]);slide.shapes.title.text='Old';slide.shapes.title.name='slot:title';box=slide.shapes.add_textbox(Inches(1),Inches(2),Inches(5),Inches(1));box.name='decorative';box.text='Never change';prs.save(self.source)
 def tearDown(self):self.temp.cleanup()
 def make_plan(self):
  summary=self.tool.inspect(self.source,view='summary');slide=self.tool.inspect(self.source,view='slide',slide_id=summary['slides'][0]['id']);target=next(x for x in slide['slots'] if x['key']=='title');return self.tool.plan({'summary':summary,'slides':[slide]},{'operations':[{'type':'set_slot_text','target_id':target['id'],'text':'New','expected_text':'Old'}]})['plan']
 def test_stale_forged_duplicate_and_source_equals_output_refuse(self):
  plan=self.make_plan();out=self.root/'out.pptx';forged=copy.deepcopy(plan);forged['operations'][0]['text']='Forged';result=self.tool.apply(self.source,forged,out);self.assertEqual(result['status'],'refused');self.assertFalse(out.exists())
  duplicate=copy.deepcopy(plan);duplicate['operations']*=2;raw=dict(duplicate);raw.pop('plan_sha256');duplicate['plan_sha256']=_obj(raw);result=self.tool.apply(self.source,duplicate,out);self.assertEqual(result['reason'],'conflict');self.assertFalse(out.exists())
  before=self.source.read_bytes();result=self.tool.apply(self.source,plan,self.source);self.assertEqual(result['reason'],'unsafe_plan');self.assertEqual(self.source.read_bytes(),before)
 def test_apply_revalidates_snapshot_expected_text_budget_and_mixed_reorder(self):
  plan=self.make_plan();out=self.root/'guarded.pptx'
  for mutate,reason in [
   (lambda p:p.update(snapshot_sha256='0'*64),'stale_snapshot'),
   (lambda p:p['operations'][0].update(expected_text='forged'),'stale_snapshot'),
   (lambda p:p['operations'][0].update(text='x'*32768),'validation_failure'),
  ]:
   forged=copy.deepcopy(plan);mutate(forged);raw=dict(forged);raw.pop('plan_sha256');forged['plan_sha256']=_obj(raw);result=self.tool.apply(self.source,forged,out);self.assertEqual(result['status'],'refused');self.assertEqual(result['reason'],reason);self.assertFalse(out.exists())
  summary=self.tool.inspect(self.source);slide=self.tool.inspect(self.source,view='slide',slide_id=summary['slides'][0]['id']);mixed={'operations':[{'type':'reorder_slides','slide_ids':[summary['slides'][0]['id']]},{'type':'set_slot_text','target_id':slide['slots'][0]['id'],'text':'x','expected_text':'Old'}]};self.assertEqual(self.tool.plan({'summary':summary,'slides':[slide]},mixed)['reason'],'conflict')

 def test_apply_refuses_forged_mixed_reorder_and_content_plan(self):
  plan=self.make_plan();summary=self.tool.inspect(self.source);ids=[slide['id'] for slide in summary['slides']];forged=copy.deepcopy(plan);forged['operations'].insert(0,{'type':'reorder_slides','slide_ids':ids,'old_slide_ids':ids});raw=dict(forged);raw.pop('plan_sha256');forged['plan_sha256']=_obj(raw);out=self.root/'mixed.pptx';result=self.tool.apply(self.source,forged,out);self.assertEqual(result['status'],'refused');self.assertEqual(result['reason'],'conflict');self.assertFalse(out.exists())

 def test_forged_apply_operation_with_extra_field_refuses(self):
  plan=self.make_plan();forged=copy.deepcopy(plan);forged['operations'][0]['unexpected']='x';raw=dict(forged);raw.pop('plan_sha256');forged['plan_sha256']=_obj(raw);out=self.root/'forged-extra.pptx';result=self.tool.apply(self.source,forged,out);self.assertEqual(result['status'],'refused');self.assertEqual(result['reason'],'validation_failure');self.assertFalse(out.exists())

 def test_candidate_collateral_change_does_not_publish(self):
  plan=self.make_plan();out=self.root/'collateral.pptx';real_publish=__import__('pptx_artifact_tool.api',fromlist=['_publish'])._publish
  def corrupting_publish(output,build,check):
   def corrupt_build(candidate):
    build(candidate)
    with zipfile.ZipFile(candidate) as src:
     payload={i.filename:src.read(i.filename) for i in src.infolist()};infos=src.infolist()
    payload['ppt/presentation.xml']=payload['ppt/presentation.xml'].replace(b'</p:presentation>',b'<!--unexpected--></p:presentation>')
    with zipfile.ZipFile(candidate,'w') as dst:
     for info in infos:dst.writestr(info,payload[info.filename])
   return real_publish(output,corrupt_build,check)
  with patch('pptx_artifact_tool.api._publish',side_effect=corrupting_publish):result=self.tool.apply(self.source,plan,out)
  self.assertEqual(result['status'],'refused');self.assertEqual(result['reason'],'validation_failure');self.assertFalse(out.exists())

 def test_candidate_postcondition_failure_does_not_publish(self):
  plan=self.make_plan();out=self.root/'never.pptx';real=Presentation
  def corrupt(path=None,*args,**kwargs):
   prs=real(path,*args,**kwargs)
   if path is not None and Path(path).parent==self.root:prs.slides[0].shapes.title.text='Wrong'
   return prs
  with patch('pptx_artifact_tool.api.Presentation',side_effect=corrupt):result=self.tool.apply(self.source,plan,out)
  self.assertEqual(result['status'],'refused');self.assertEqual(result['reason'],'validation_failure');self.assertFalse(out.exists())

 def test_apply_uses_private_source_snapshot_after_preflight(self):
  plan=self.make_plan();out=self.root/'snapshot-bound.pptx';replacement=self.root/'replacement.pptx';prs=Presentation();slide=prs.slides.add_slide(prs.slide_layouts[5]);slide.shapes.title.text='Old';slide.shapes.title.name='slot:title';box=slide.shapes.add_textbox(Inches(1),Inches(2),Inches(5),Inches(1));box.name='decorative';box.text='Injected collateral';prs.save(replacement)
  real_publish=__import__('pptx_artifact_tool.api',fromlist=['_publish'])._publish
  def swapping_publish(output,build,check):
   self.source.write_bytes(replacement.read_bytes())
   return real_publish(output,build,check)
  with patch('pptx_artifact_tool.api._publish',side_effect=swapping_publish):result=self.tool.apply(self.source,plan,out)
  self.assertEqual(result['status'],'ok');rendered=Presentation(out);self.assertEqual(rendered.slides[0].shapes.title.text,'New');self.assertEqual(next(shape.text for shape in rendered.slides[0].shapes if shape.name=='decorative'),'Never change')

 def test_inspect_refuses_decks_over_slide_budget(self):
  oversized=self.root/'oversized.pptx';prs=Presentation()
  for _ in range(501):prs.slides.add_slide(prs.slide_layouts[6])
  prs.save(oversized);result=self.tool.inspect(oversized)
  self.assertEqual(result['status'],'refused');self.assertEqual(result['reason'],'unsafe_plan')

 def test_inspect_enforces_aggregate_slot_budget(self):
  deck=self.root/'slot-budget.pptx';prs=Presentation();slide=prs.slides.add_slide(prs.slide_layouts[6])
  for index in range(2):shape=slide.shapes.add_textbox(Inches(1),Inches(1+index),Inches(2),Inches(0.5));shape.name=f'slot:text-{index}';shape.text=str(index)
  prs.save(deck)
  with patch('pptx_artifact_tool.api.MAX_SLOTS',1):result=self.tool.inspect(deck)
  self.assertEqual(result['status'],'refused');self.assertEqual(result['reason'],'unsafe_plan')

 def test_inspect_enforces_aggregate_table_cell_budget(self):
  deck=self.root/'cell-budget.pptx';prs=Presentation();slide=prs.slides.add_slide(prs.slide_layouts[6]);shape=slide.shapes.add_table(1,2,Inches(1),Inches(1),Inches(4),Inches(1));shape.name='slot:data';prs.save(deck)
  with patch('pptx_artifact_tool.api.MAX_TABLE_CELLS',1):result=self.tool.inspect(deck)
  self.assertEqual(result['status'],'refused');self.assertEqual(result['reason'],'unsafe_plan')

 def test_plan_refuses_forged_snapshot_over_cardinality_budgets(self):
  summary=self.tool.inspect(self.source,view='summary');slide=self.tool.inspect(self.source,view='slide',slide_id=summary['slides'][0]['id']);request={'operations':[{'type':'set_slot_text','target_id':slide['slots'][0]['id'],'text':'New','expected_text':'Old'}]}
  cases=(
   ({'summary':summary,'slides':[slide,copy.deepcopy(slide)]},'MAX_SLIDES'),
   ({'summary':summary,'slides':[{**slide,'slots':slide['slots']+[copy.deepcopy(slide['slots'][0])]}]},'MAX_SLOTS'),
   ({'summary':summary,'slides':[{**slide,'table_cells':[{'id':'cell','text':'','row':0,'column':0},{'id':'cell-2','text':'','row':0,'column':1}]}]},'MAX_TABLE_CELLS'),
  )
  for snapshot,limit_name in cases:
   with self.subTest(limit=limit_name),patch(f'pptx_artifact_tool.api.{limit_name}',1):
    result=self.tool.plan(snapshot,request)
    self.assertEqual(result['status'],'refused');self.assertEqual(result['reason'],'unsafe_plan')

 def test_unnamed_shape_and_unsafe_zip_are_not_mutable(self):
  summary=self.tool.inspect(self.source,view='summary');slide=self.tool.inspect(self.source,view='slide',slide_id=summary['slides'][0]['id']);self.assertNotIn('decorative',[x['key'] for x in slide['slots']]);result=self.tool.plan({'summary':summary,'slides':[slide]},{'operations':[{'type':'set_slot_text','target_id':'tx_000000000000000000000000','text':'x','expected_text':'Never change'}]});self.assertEqual(result['status'],'refused')
  unsafe=self.root/'unsafe.pptx'
  with zipfile.ZipFile(unsafe,'w') as z:z.writestr('../escape.xml','x')
  self.assertEqual(self.tool.inspect(unsafe)['status'],'refused');self.assertEqual(self.tool.validate(unsafe)['status'],'invalid')

if __name__=='__main__':unittest.main()
