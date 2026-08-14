from __future__ import annotations
import copy,hashlib,json,os,stat,tempfile,zipfile
from pathlib import Path,PurePosixPath
from typing import Any
from lxml import etree
from pptx import Presentation

PML='http://schemas.openxmlformats.org/presentationml/2006/main';DML='http://schemas.openxmlformats.org/drawingml/2006/main';REL='http://schemas.openxmlformats.org/package/2006/relationships';OFFREL='http://schemas.openxmlformats.org/officeDocument/2006/relationships';NS={'p':PML,'a':DML,'r':OFFREL}
MAX_SLIDES=500;MAX_SLOTS=5000;MAX_TABLE_CELLS=20000
def _sha(path:Path):return hashlib.sha256(path.read_bytes()).hexdigest()
def _obj(value):return hashlib.sha256(json.dumps(value,sort_keys=True,separators=(',',':'),ensure_ascii=False).encode()).hexdigest()
def _refuse(reason,details=''):return {'status':'refused','reason':reason,'details':details}
def _id(*parts):return 'tx_'+hashlib.sha256('|'.join(map(str,parts)).encode()).hexdigest()[:24]
def _admit(path:Path):
 try:
  if path.suffix.lower()!='.pptx':raise ValueError
  with zipfile.ZipFile(path) as z:
   infos=z.infolist();names=[x.filename for x in infos];total=sum(x.file_size for x in infos);modes=[(x.external_attr>>16)&0xFFFF for x in infos]
   unsafe=not {'[Content_Types].xml','ppt/presentation.xml'}<=set(names) or len(names)!=len(set(names)) or any(n.startswith('/') or '\\' in n or '..' in PurePosixPath(n).parts for n in names) or any(not(i.is_dir() or stat.S_IFMT(m) in (0,stat.S_IFREG)) for i,m in zip(infos,modes)) or len(infos)>10000 or total>128*1024*1024 or any(x.file_size>32*1024*1024 for x in infos)
   if unsafe:raise ValueError
 except Exception as exc:raise ValueError('unsafe_package') from exc
def _source_snapshot(source:Path,workdir:Path):
 source_fd=-1;snapshot=None
 try:
  source_fd=os.open(source,os.O_RDONLY|getattr(os,'O_NOFOLLOW',0))
  if not stat.S_ISREG(os.fstat(source_fd).st_mode):raise ValueError('validation_failure')
  fd,name=tempfile.mkstemp(prefix='.source-snapshot.',suffix='.pptx',dir=workdir);snapshot=Path(name)
  with os.fdopen(source_fd,'rb') as original,os.fdopen(fd,'wb') as target:
   source_fd=-1
   for block in iter(lambda:original.read(1024*1024),b''):target.write(block)
  return snapshot
 except Exception as exc:
  if source_fd>=0:os.close(source_fd)
  if snapshot is not None:snapshot.unlink(missing_ok=True)
  raise ValueError('validation_failure') from exc
def _publish(output:Path,build,check):
 output.parent.mkdir(parents=True,exist_ok=True);fd,name=tempfile.mkstemp(prefix='.'+output.name+'.candidate.',suffix='.pptx',dir=output.parent);os.close(fd);candidate=Path(name)
 try:build(candidate);check(candidate);os.replace(candidate,output)
 finally:candidate.unlink(missing_ok=True)
def _shape_text(shape):return shape.text if getattr(shape,'has_text_frame',False) else ''
def _slot(shape):return shape.name[5:] if shape.name.startswith('slot:') else None
def _slide_members(source:Path):
 with zipfile.ZipFile(source) as archive:
  root=etree.fromstring(archive.read('ppt/presentation.xml'));rels=etree.fromstring(archive.read('ppt/_rels/presentation.xml.rels'));targets={node.get('Id'):node.get('Target') for node in rels.findall(f'{{{REL}}}Relationship')}
 ids=root.xpath('./p:sldIdLst/p:sldId/@r:id',namespaces=NS);members=[]
 for rid in ids:
  target=targets.get(rid)
  if not target:raise ValueError('validation_failure')
  if target.startswith('/'):member=target.lstrip('/')
  elif target.startswith('ppt/'):member=target
  else:member=str(PurePosixPath('ppt')/target)
  normalized=str(PurePosixPath(member))
  if normalized not in archive.namelist():raise ValueError('validation_failure')
  members.append(normalized)
 return members
def _slide_member(source:Path,index):
 members=_slide_members(source)
 if not isinstance(index,int) or index<0 or index>=len(members):raise ValueError('validation_failure')
 return members[index]
def _slide_payload(source:Path,index:int,slide):
 digest=_sha(source);slots=[];cells=[]
 for shape_index,shape in enumerate(slide.shapes):
  key=_slot(shape)
  if key and shape.has_text_frame:
   text=_shape_text(shape);slots.append({'id':_id(digest,index,shape.shape_id,'slot',text),'key':key,'kind':'text','shape_name':shape.name,'text':text,'shape_id':shape.shape_id})
  if key and shape.has_table:
   slots.append({'id':_id(digest,index,shape.shape_id,'table'),'key':key,'kind':'table','shape_name':shape.name,'text':'','shape_id':shape.shape_id})
   for r,row in enumerate(shape.table.rows):
    for c,cell in enumerate(row.cells):
     cells.append({'id':_id(digest,index,shape.shape_id,r,c,cell.text),'slot_key':key,'key':f'{key}.r{r}c{c}','row':r,'column':c,'text':cell.text,'shape_name':shape.name,'shape_id':shape.shape_id})
 return {'id':_id(digest,'slide',index),'index':index,'title':slide.shapes.title.text if slide.shapes.title else '', 'layout':slide.slide_layout.name,'slots':slots,'table_cells':cells}
def _summary(source:Path):
 _admit(source);prs=Presentation(source)
 if len(prs.slides)>MAX_SLIDES:raise ValueError('unsafe_plan')
 slides=[_slide_payload(source,i,s) for i,s in enumerate(prs.slides)]
 if sum(len(s['slots']) for s in slides)>MAX_SLOTS or sum(len(s['table_cells']) for s in slides)>MAX_TABLE_CELLS:raise ValueError('unsafe_plan')
 summary={'status':'ok','artifact_type':'pptx','view':'summary','source_sha256':_sha(source),'slides':[{k:v for k,v in s.items() if k!='table_cells'} for s in slides]};summary['snapshot_sha256']=_obj({k:v for k,v in summary.items() if k!='status'});return summary
def _set_text_nodes(shape_element,text):
 nodes=shape_element.xpath('.//a:t',namespaces=NS)
 if not nodes:raise ValueError('unsupported_capability')
 first=nodes[0];run=first.getparent();paragraph=run.getparent();parts=text.split('\n');first.text=parts[0]
 for node in nodes[1:]:node.text=''
 insert_at=paragraph.index(run)+1
 for part in parts[1:]:
  paragraph.insert(insert_at,etree.Element(f'{{{DML}}}br'));insert_at+=1
  new_run=copy.deepcopy(run);new_text=new_run.xpath('./a:t',namespaces=NS)
  if len(new_text)!=1:raise ValueError('unsupported_capability')
  new_text[0].text=part;paragraph.insert(insert_at,new_run);insert_at+=1
def _mutate(source:Path,candidate:Path,ops):
 replacements={};by_slide={};reorder=None
 for op in ops:
  if op['type']=='reorder_slides':reorder=op;continue
  by_slide.setdefault(op['slide_index'],[]).append(op)
 with zipfile.ZipFile(source) as archive:
  if reorder:
   root=etree.fromstring(archive.read('ppt/presentation.xml'));lst=root.xpath('./p:sldIdLst',namespaces=NS)
   if len(lst)!=1:raise ValueError('validation_failure')
   current=list(lst[0]);mapping={ident:node for ident,node in zip(reorder['old_slide_ids'],current)}
   for node in current:lst[0].remove(node)
   for ident in reorder['slide_ids']:lst[0].append(mapping[ident])
   replacements['ppt/presentation.xml']=etree.tostring(root,xml_declaration=True,encoding='UTF-8',standalone=True)
  for slide_index,slide_ops in by_slide.items():
   member=_slide_member(source,slide_index);root=etree.fromstring(archive.read(member));shapes=root.xpath('.//p:sp|.//p:graphicFrame',namespaces=NS)
   for op in slide_ops:
    matched=[shape for shape in shapes if shape.xpath('./p:nvSpPr/p:cNvPr[@id=$shape_id] | ./p:nvGraphicFramePr/p:cNvPr[@id=$shape_id]',namespaces=NS,shape_id=str(op['shape_id']))]
    if len(matched)!=1:raise ValueError('ambiguous_target')
    shape_element=matched[0]
    if op['type'] in {'set_slot_text','clear_slot_text'}:_set_text_nodes(shape_element,op.get('text',''))
    else:
     table=shape_element.xpath('.//a:tbl',namespaces=NS)
     if len(table)!=1:raise ValueError('validation_failure')
     rows=table[0].xpath('./a:tr',namespaces=NS);cells=rows[op['row']].xpath('./a:tc',namespaces=NS);_set_text_nodes(cells[op['column']],op['text'])
   replacements[member]=etree.tostring(root,xml_declaration=True,encoding='UTF-8',standalone=True)
  with zipfile.ZipFile(candidate,'w') as output:
   for info in archive.infolist():output.writestr(copy.copy(info),replacements.get(info.filename,archive.read(info.filename)))

class PptxArtifactTool:
 def __init__(self,workdir:Path|str):self.workdir=Path(workdir);self.workdir.mkdir(parents=True,exist_ok=True)
 def inspect(self,source:Path|str,view='summary',slide_id=None,query=None):
  try:
   source=Path(source);summary=_summary(source)
   if view=='summary':return summary
   prs=Presentation(source);slides=[_slide_payload(source,i,s) for i,s in enumerate(prs.slides)]
   if view=='slide':
    found=[s for s in slides if s['id']==slide_id]
    if len(found)!=1:raise ValueError('ambiguous_target')
    result={**found[0],'status':'ok','artifact_type':'pptx','view':'slide','source_sha256':_sha(source),'snapshot_sha256':summary['snapshot_sha256']};return result
   if view=='search':
    if not isinstance(query,str) or not query:raise ValueError('validation_failure')
    matches=[]
    for s in slides:
     for item in s['slots']+s['table_cells']:
      if query.casefold() in item['text'].casefold():matches.append({'slide_id':s['id'],'slide_index':s['index'],**item})
    return {'status':'ok','artifact_type':'pptx','view':'search','source_sha256':_sha(source),'query':query,'matches':matches[:100],'truncated':len(matches)>100}
   raise ValueError('unsupported_view')
  except (ValueError,KeyError,TypeError) as exc:return _refuse(str(exc))
 def plan(self,snapshot:dict[str,Any],request:dict[str,Any]):
  try:
   if not isinstance(snapshot,dict) or not isinstance(request,dict) or set(request)!={'operations'}:raise ValueError('validation_failure')
   summary=snapshot.get('summary');slides=snapshot.get('slides')
   if not isinstance(summary,dict) or summary.get('status')!='ok' or not isinstance(slides,list):raise ValueError('validation_failure')
   if len(slides)>MAX_SLIDES:raise ValueError('unsafe_plan')
   total_slots=0;total_table_cells=0
   for slide in slides:
    if not isinstance(slide,dict) or not isinstance(slide.get('slots'),list) or not isinstance(slide.get('table_cells'),list):raise ValueError('validation_failure')
    total_slots+=len(slide['slots']);total_table_cells+=len(slide['table_cells'])
    if total_slots>MAX_SLOTS or total_table_cells>MAX_TABLE_CELLS:raise ValueError('unsafe_plan')
   targets={}
   for slide in slides:
    for item in slide.get('slots',[]):targets[item['id']]={**item,'slide_index':slide['index']}
    for item in slide.get('table_cells',[]):targets[item['id']]={**item,'slide_index':slide['index']}
   ops=[];seen=set()
   slide_by_id={slide['id']:slide for slide in slides}
   for op in request['operations']:
    if isinstance(op,dict) and op.get('type')=='reorder_slides':
     if set(op)!={'type','slide_ids'} or not isinstance(op['slide_ids'],list) or len(op['slide_ids'])!=len(slide_by_id) or set(op['slide_ids'])!=set(slide_by_id):raise ValueError('validation_failure')
     ops.append({'type':'reorder_slides','slide_ids':op['slide_ids'],'old_slide_ids':[slide['id'] for slide in sorted(slides,key=lambda x:x['index'])]});continue
    if not isinstance(op,dict) or not isinstance(op.get('target_id'),str) or op['target_id'] not in targets:raise ValueError('ambiguous_target')
    if op['target_id'] in seen:raise ValueError('conflict')
    seen.add(op['target_id']);target=targets[op['target_id']];kind=op.get('type')
    if kind in {'set_slot_text','clear_slot_text'}:
     allowed={'type','target_id','expected_text'}|({'text'} if kind=='set_slot_text' else set())
     if set(op)!=allowed or target.get('kind')!='text' or op['expected_text']!=target['text'] or kind=='set_slot_text' and (not isinstance(op['text'],str) or len(op['text'])>32767):raise ValueError('validation_failure')
     resolved={**op,**{k:target[k] for k in ('slide_index','shape_id','shape_name')},'old_text':target['text']}
    elif kind=='set_table_cell_text':
     if set(op)!={'type','target_id','text','expected_text'} or 'row' not in target or op['expected_text']!=target['text'] or not isinstance(op['text'],str) or len(op['text'])>32767:raise ValueError('validation_failure')
     resolved={**op,**{k:target[k] for k in ('slide_index','shape_id','shape_name','row','column')},'old_text':target['text']}
    else:raise ValueError('unsupported_capability')
    ops.append(resolved)
   if any(op['type']=='reorder_slides' for op in ops) and len(ops)>1:raise ValueError('conflict')
   if not ops or len(ops)>1000:raise ValueError('unsafe_plan')
   plan={'schema':1,'source_sha256':summary['source_sha256'],'snapshot_sha256':summary['snapshot_sha256'],'operations':ops};plan['plan_sha256']=_obj(plan);return {'status':'ok','plan':plan}
  except (ValueError,KeyError,TypeError) as exc:return _refuse(str(exc))
 def apply(self,source:Path|str,plan:dict[str,Any],output:Path|str):
  source=Path(source);output=Path(output);source_snapshot=None
  try:
   if source.resolve()==output.resolve():raise ValueError('unsafe_plan')
   if not isinstance(plan,dict) or set(plan)!={'schema','source_sha256','snapshot_sha256','operations','plan_sha256'} or not isinstance(plan['operations'],list) or not plan['operations'] or len(plan['operations'])>1000:raise ValueError('validation_failure')
   check=dict(plan);provided=check.pop('plan_sha256')
   if provided!=_obj(check):raise ValueError('validation_failure')
   source_snapshot=_source_snapshot(source,self.workdir);source=source_snapshot;_admit(source)
   if plan['source_sha256']!=_sha(source):raise ValueError('stale_snapshot')
   current=_summary(source)
   if plan['snapshot_sha256']!=current['snapshot_sha256']:raise ValueError('stale_snapshot')
   prs=Presentation(source);targets={}
   for i,slide in enumerate(prs.slides):
    payload=_slide_payload(source,i,slide)
    for item in payload['slots']+payload['table_cells']:targets[item['id']]={**item,'slide_index':i}
   seen=set();current_slide_ids=[_id(_sha(source),'slide',i) for i in range(len(prs.slides))]
   allowed_fields={
    'set_slot_text':{'type','target_id','text','expected_text','slide_index','shape_id','shape_name','old_text'},
    'clear_slot_text':{'type','target_id','expected_text','slide_index','shape_id','shape_name','old_text'},
    'set_table_cell_text':{'type','target_id','text','expected_text','slide_index','shape_id','shape_name','row','column','old_text'},
    'reorder_slides':{'type','slide_ids','old_slide_ids'},
   }
   for op in plan['operations']:
    if not isinstance(op,dict) or op.get('type') not in allowed_fields or set(op)!=allowed_fields[op['type']]:raise ValueError('validation_failure')
    if op['type']=='reorder_slides':
     if set(op)!={'type','slide_ids','old_slide_ids'} or op['old_slide_ids']!=current_slide_ids or len(op['slide_ids'])!=len(current_slide_ids) or set(op['slide_ids'])!=set(current_slide_ids):raise ValueError('stale_snapshot')
     continue
    if not isinstance(op,dict) or op.get('target_id') in seen or op.get('target_id') not in targets:raise ValueError('conflict' if op.get('target_id') in seen else 'stale_snapshot')
    seen.add(op['target_id']);target=targets[op['target_id']]
    if op.get('old_text')!=target['text'] or op.get('expected_text')!=target['text'] or op.get('slide_index')!=target['slide_index'] or op.get('shape_id')!=target['shape_id'] or op.get('shape_name')!=target['shape_name']:raise ValueError('stale_snapshot')
    if op['type']=='set_slot_text' and (not isinstance(op.get('text'),str) or len(op['text'])>32767):raise ValueError('validation_failure')
    if op['type']=='set_table_cell_text':
     if not isinstance(op.get('text'),str) or len(op['text'])>32767:raise ValueError('validation_failure')
     if op.get('row')!=target.get('row') or op.get('column')!=target.get('column'):raise ValueError('stale_snapshot')
   if any(op['type']=='reorder_slides' for op in plan['operations']) and len(plan['operations'])>1:raise ValueError('conflict')
   def check_candidate(candidate):
    _admit(candidate)
    allowed_members={_slide_member(source,op['slide_index']) for op in plan['operations'] if op['type']!='reorder_slides'}
    if any(op['type']=='reorder_slides' for op in plan['operations']):allowed_members.add('ppt/presentation.xml')
    with zipfile.ZipFile(source) as original,zipfile.ZipFile(candidate) as changed:
     original_hashes={i.filename:hashlib.sha256(original.read(i.filename)).hexdigest() for i in original.infolist()};changed_hashes={i.filename:hashlib.sha256(changed.read(i.filename)).hexdigest() for i in changed.infolist()}
    actual_changed={name for name in original_hashes.keys()|changed_hashes.keys() if original_hashes.get(name)!=changed_hashes.get(name)}
    if set(original_hashes)!=set(changed_hashes) or not actual_changed.issubset(allowed_members):raise ValueError('validation_failure')
    after=Presentation(candidate)
    for op in plan['operations']:
     if op['type']=='reorder_slides':
      titles=[slide.shapes.title.text if slide.shapes.title else '' for slide in after.slides];before_titles=[slide.shapes.title.text if slide.shapes.title else '' for slide in prs.slides];mapping=dict(zip(current_slide_ids,before_titles))
      if titles!=[mapping[ident] for ident in op['slide_ids']]:raise ValueError('validation_failure')
      continue
     candidates=[shape for shape in after.slides[op['slide_index']].shapes if shape.shape_id==op['shape_id']]
     if len(candidates)!=1:raise ValueError('validation_failure')
     shape=candidates[0]
     actual=shape.text if op['type']!='set_table_cell_text' else shape.table.cell(op['row'],op['column']).text
     expected=op.get('text','').replace('\n','\v')
     if actual!=expected:raise ValueError('validation_failure')
   _publish(output,lambda p:_mutate(source,p,plan['operations']),check_candidate);return {'status':'ok','sha256':_sha(output),'changed_targets':[op['target_id'] for op in plan['operations'] if 'target_id' in op],'slide_order_changed':any(op['type']=='reorder_slides' for op in plan['operations'])}
  except (ValueError,KeyError,TypeError,IndexError,zipfile.BadZipFile) as exc:return _refuse(str(exc))
  finally:
   if source_snapshot is not None:source_snapshot.unlink(missing_ok=True)
 def validate(self,source:Path|str,before:Path|str|None=None):
  source=Path(source)
  try:
   _admit(source);Presentation(source);changed=[]
   if before:
    with zipfile.ZipFile(before) as a,zipfile.ZipFile(source) as b:
     ah={i.filename:hashlib.sha256(a.read(i.filename)).hexdigest() for i in a.infolist()};bh={i.filename:hashlib.sha256(b.read(i.filename)).hexdigest() for i in b.infolist()};changed=sorted(k for k in ah.keys()|bh.keys() if ah.get(k)!=bh.get(k))
   unexpected=[x for x in changed if not (x.startswith('ppt/slides/slide') or x=='ppt/presentation.xml')]
   return {'status':'valid' if not unexpected else 'invalid','package_valid':True,'changed_members':changed,'unexpected_changed_members':unexpected,'application_compatibility':'not_executed','visual_fidelity':'not_executed'}
  except Exception:return {'status':'invalid','error':'invalid_package'}
 def create(self,template:Path|str,model:dict[str,Any],output:Path|str):
  template=Path(template);summary=self.inspect(template,view='summary')
  if summary.get('status')!='ok' or not isinstance(model,dict) or set(model)-{'slots','table_cells'}:return _refuse('validation_failure')
  slides=[self.inspect(template,view='slide',slide_id=s['id']) for s in summary['slides']];all_slots=[x for s in slides for x in s['slots'] if x.get('kind')=='text'];all_cells=[x for s in slides for x in s['table_cells']]
  slot_keys=[x['key'] for x in all_slots];cell_keys=[x['key'] for x in all_cells]
  if len(slot_keys)!=len(set(slot_keys)) or len(cell_keys)!=len(set(cell_keys)):return _refuse('ambiguous_target')
  slot_map={x['key']:x for x in all_slots};cell_map={x['key']:x for x in all_cells};ops=[]
  try:
   for key,text in model.get('slots',{}).items():
    target=slot_map[key];ops.append({'type':'set_slot_text','target_id':target['id'],'text':text,'expected_text':target['text']})
   for key,text in model.get('table_cells',{}).items():
    target=cell_map[key];ops.append({'type':'set_table_cell_text','target_id':target['id'],'text':text,'expected_text':target['text']})
  except (KeyError,TypeError):return _refuse('ambiguous_target')
  planned=self.plan({'summary':summary,'slides':slides},{'operations':ops})
  return planned if planned.get('status')!='ok' else self.apply(template,planned['plan'],output)
