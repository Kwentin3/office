"""Independent bounded validation for generated PPTX candidates.

STICKY CLAIM BOUNDARY: parser/package/geometry checks are not visual rendering
and are not PowerPoint compatibility. Gates that did not run remain explicit
`not_executed`; never collapse them into PASS merely because python-pptx opens.
"""

from __future__ import annotations

import shutil
import zipfile
from pathlib import Path, PurePosixPath
from typing import Any

from lxml import etree
from pptx import Presentation

from .contracts import validate_deck_spec


_REQUIRED_PARTS = {"[Content_Types].xml", "_rels/.rels", "ppt/presentation.xml"}


def _issue(code: str, message: str, **extra: Any) -> dict[str, Any]:
    return {"code": code, "message": message, **extra}


def _package_checks(path: Path) -> list[dict[str, Any]]:
    issues: list[dict[str, Any]] = []
    try:
        with zipfile.ZipFile(path) as package:
            infos = package.infolist()
            names = [info.filename for info in infos]
            if len(names) != len(set(names)):
                issues.append(_issue("duplicate_zip_member", "PPTX contains duplicate ZIP member names"))
            if len(infos) > 10000:
                issues.append(_issue("package_member_budget", "PPTX exceeds package member budget"))
            total = 0
            for info in infos:
                pure = PurePosixPath(info.filename)
                if info.filename.startswith(("/", "\\")) or "\\" in info.filename or ".." in pure.parts:
                    issues.append(_issue("unsafe_zip_member", f"Unsafe member: {info.filename}"))
                if info.file_size > 32 * 1024 * 1024:
                    issues.append(_issue("package_member_size", f"Member too large: {info.filename}"))
                total += info.file_size
                if info.filename.endswith((".xml", ".rels")):
                    try:
                        etree.fromstring(package.read(info))
                    except Exception as exc:
                        issues.append(_issue("malformed_xml", f"Malformed XML: {info.filename}: {exc}"))
            if total > 128 * 1024 * 1024:
                issues.append(_issue("package_size", "PPTX exceeds uncompressed package budget"))
            for required in sorted(_REQUIRED_PARTS - set(names)):
                issues.append(_issue("missing_required_part", f"Missing required package part: {required}"))
    except (OSError, zipfile.BadZipFile) as exc:
        issues.append(_issue("invalid_zip", str(exc)))
    return issues


def _source_citation(slide: dict[str, Any], source_labels: dict[str, str]) -> str | None:
    labels = [source_labels[source_id] for source_id in slide["source_ids"]]
    if not labels:
        return None
    prefix = "Источник: " if len(labels) == 1 else "Источники: "
    return prefix + "; ".join(labels)
def _expected_texts(slide: dict[str, Any], source_labels: dict[str, str]) -> list[str]:
    archetype = slide["archetype"]
    if archetype == "cover":
        expected = [slide["title"], slide["subtitle"]]
    elif archetype == "comparison":
        expected = [slide["title"], slide["left"]["label"], slide["right"]["label"]]
        expected.extend(slide["left"]["items"]); expected.extend(slide["right"]["items"])
    elif archetype == "chart_with_takeaway":
        expected = [slide["title"], slide["takeaway"]]
    elif archetype == "process":
        expected = [slide["title"]] + [value for step in slide["steps"] for value in (step["label"], step["description"])]
    elif archetype == "timeline":
        expected = [slide["title"]] + [value for milestone in slide["milestones"] for value in (milestone["period"], milestone["label"])]
    elif archetype == "decision_matrix":
        expected = [slide["title"]] + slide["criteria"] + [option["label"] for option in slide["options"]]
    else:
        expected = [slide["title"]] + [value for metric in slide["metrics"] for value in (metric["label"], metric["value"], metric["note"])]
    citation = _source_citation(slide, source_labels)
    if citation:
        expected.append(citation)
    return expected


def _parser_semantic_geometry_checks(path: Path, deck: dict[str, Any]) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    semantic: list[dict[str, Any]] = []
    geometry: list[dict[str, Any]] = []
    try:
        prs = Presentation(path)
    except Exception as exc:
        return [_issue("parser_open_failed", str(exc))], []
    if len(prs.slides) != len(deck["slides"]):
        semantic.append(_issue("slide_count_mismatch", f"Expected {len(deck['slides'])}, got {len(prs.slides)}"))
        return semantic, geometry
    source_labels = {source["source_id"]: source["label"] for source in deck["sources"]}
    for index, (slide, expected) in enumerate(zip(prs.slides, deck["slides"])):
        text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        citation = _source_citation(expected, source_labels)
        for expected_text in _expected_texts(expected, source_labels):
            if expected_text not in text:
                code = "missing_source_citation" if expected_text == citation else "missing_expected_text"
                semantic.append(_issue(code, f"Expected text not found: {expected_text}", slide_id=expected["slide_id"]))
        expected_chart = expected["archetype"] == "chart_with_takeaway"
        if expected_chart and not any(shape.has_chart for shape in slide.shapes):
            semantic.append(_issue("missing_native_chart", "Expected a native chart", slide_id=expected["slide_id"]))
        for shape in slide.shapes:
            # STICKY LIMIT: bounding-box checks catch off-canvas defects, not text
            # overflow or rotated-shape collisions. Those require a real renderer.
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
                geometry.append(_issue("off_slide_shape", f"Shape {shape.name!r} crosses the slide boundary", slide_id=expected["slide_id"]))
    return semantic, geometry


def validate_presentation(source: str | Path, deck_spec: dict[str, Any]) -> dict[str, Any]:
    """Validate the executable subset and report unavailable gates honestly."""
    deck = validate_deck_spec(deck_spec)
    path = Path(source)
    structural_issues = _package_checks(path)
    semantic_issues: list[dict[str, Any]] = []
    geometry_issues: list[dict[str, Any]] = []
    if not structural_issues:
        semantic_issues, geometry_issues = _parser_semantic_geometry_checks(path, deck)
    invalid = bool(structural_issues or semantic_issues or geometry_issues)
    visual_renderer = shutil.which("soffice") or shutil.which("libreoffice")
    if visual_renderer:
        visual = {"status": "not_executed", "reason": "renderer discovered but raster pipeline is not implemented in this vertical slice"}
    else:
        visual = {"status": "not_executed", "reason": "visual renderer unavailable; overflow/collision/pixel checks were not executed"}
    application = {"status": "not_executed", "reason": "PowerPoint application open/save gate is unavailable in this environment"}
    return {
        "status": "invalid" if invalid else "valid_with_unexecuted_gates",
        "structural": {"status": "fail" if structural_issues else "pass", "issues": structural_issues},
        "semantic": {"status": "fail" if semantic_issues else "pass", "issues": semantic_issues},
        "geometry": {"status": "fail" if geometry_issues else "pass", "issues": geometry_issues},
        "visual": visual,
        "application": application,
    }
