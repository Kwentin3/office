"""Native editable PPTX backend for validated bounded SceneSpec.

STICKY DOMAIN BOUNDARY: this backend knows finite scene primitives, PPTX object
construction, and no semantic archetype recipes. Asset admission and DeckSpec
validation happen in the orchestration wrapper before SceneSpec reaches here.
"""

from __future__ import annotations

import copy
import hashlib
import io
import os
import tempfile
from pathlib import Path
from typing import Any, Iterable

from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from .asset_admission import AdmittedAsset, AssetAdmissionError, admit_assets
from .compiler import compile_deck
from .contracts import validate_deck_spec
from .scene_contract import validate_scene_spec


class RenderError(ValueError):
    """A deterministic refusal before final publication."""


SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.upper())


def _asset_map(deck: dict[str, Any]) -> dict[str, AdmittedAsset]:
    try:
        return admit_assets(deck)
    except AssetAdmissionError as exc:
        raise RenderError(str(exc)) from exc


def _emu_box(box: dict[str, float], canvas: dict[str, float]) -> tuple[Emu, Emu, Emu, Emu]:
    x = Emu(round(SLIDE_W * box["x"] / canvas["width"]))
    y = Emu(round(SLIDE_H * box["y"] / canvas["height"]))
    w = Emu(round(SLIDE_W * box["w"] / canvas["width"]))
    h = Emu(round(SLIDE_H * box["h"] / canvas["height"]))
    return x, y, w, h


def _shape_type(kind: str) -> MSO_AUTO_SHAPE_TYPE:
    return {"rect": MSO_AUTO_SHAPE_TYPE.RECTANGLE, "round_rect": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL}[kind]


def _add_shape(slide: Any, node: dict[str, Any], canvas: dict[str, float]) -> Any:
    x, y, w, h = _emu_box(node["box"], canvas)
    if node["shape"] == "line":
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.LINE_INVERSE, x, y, w, h)
        shape.line.color.rgb = _rgb(node.get("stroke", node["fill"]))
    else:
        shape = slide.shapes.add_shape(_shape_type(node["shape"]), x, y, w, h)
        shape.fill.solid(); shape.fill.fore_color.rgb = _rgb(node["fill"])
        if "opacity" in node:
            shape.fill.transparency = round((1.0 - node["opacity"]) * 100)
        if "stroke" in node:
            shape.line.color.rgb = _rgb(node["stroke"])
        else:
            shape.line.fill.background()
    shape.name = f"scene:{node['role']}"
    return shape


def _add_text(slide: Any, node: dict[str, Any], canvas: dict[str, float]) -> Any:
    x, y, w, h = _emu_box(node["box"], canvas)
    shape = slide.shapes.add_textbox(x, y, w, h)
    shape.name = f"scene:{node['role']}"
    frame = shape.text_frame
    frame.clear(); frame.word_wrap = True
    frame.margin_left = frame.margin_right = Emu(0)
    frame.margin_top = frame.margin_bottom = Emu(0)
    style = node["style"]
    frame.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[style["valign"]]
    paragraph = frame.paragraphs[0]
    paragraph.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[style["align"]]
    run = paragraph.add_run(); run.text = node["text"]
    run.font.name = style["font"]; run.font.size = Pt(style["size"]); run.font.bold = style["bold"]; run.font.color.rgb = _rgb(style["color"])
    return shape


def _add_chart(slide: Any, node: dict[str, Any], canvas: dict[str, float]) -> Any:
    chart_spec = node["chart"]
    data = ChartData(); data.categories = chart_spec["categories"]
    for series in chart_spec["series"]:
        data.add_series(series["name"], series["values"])
    chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED if chart_spec["type"] == "bar" else XL_CHART_TYPE.LINE_MARKERS
    x, y, w, h = _emu_box(node["box"], canvas)
    shape = slide.shapes.add_chart(chart_type, x, y, w, h, data)
    shape.name = f"scene:{node['role']}"
    chart = shape.chart
    chart.has_legend = len(chart_spec["series"]) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM; chart.legend.include_in_layout = False
    chart.has_title = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = _rgb("D9DEE3")
    palette = ["D95D39", "2A7F72", "62707D"]
    for index, series in enumerate(chart.series):
        series.format.fill.solid(); series.format.fill.fore_color.rgb = _rgb(palette[index]); series.format.line.color.rgb = _rgb(palette[index])
    return shape


def _add_image(
    slide: Any,
    node: dict[str, Any],
    canvas: dict[str, float],
    assets: dict[str, AdmittedAsset],
) -> Any:
    if node["asset_id"] not in assets:
        raise RenderError(f"scene references unavailable asset: {node['asset_id']}")
    raster_bytes = assets[node["asset_id"]].raster_bytes
    x, y, w, h = _emu_box(node["box"], canvas)
    if node["fit"] == "contain":
        picture = slide.shapes.add_picture(io.BytesIO(raster_bytes), x, y, w, h)
    else:
        with Image.open(io.BytesIO(raster_bytes)) as image:
            source_w, source_h = image.size
            target_ratio = w / h
            source_ratio = source_w / source_h
            if source_ratio >= target_ratio:
                crop_h = source_h
                crop_w = round(source_h * target_ratio)
            else:
                crop_w = source_w
                crop_h = round(source_w / target_ratio)
            left = max(0, (source_w - crop_w) // 2)
            top = max(0, (source_h - crop_h) // 2)
            crop = image.crop((left, top, left + crop_w, top + crop_h))
            snapshot = io.BytesIO()
            crop.save(snapshot, "PNG")
            snapshot.seek(0)
            picture = slide.shapes.add_picture(snapshot, x, y, w, h)
    picture.name = f"scene:{node['role']}"
    return picture


def render_scene_presentation(
    scene_spec: dict[str, Any],
    *,
    assets: dict[str, AdmittedAsset],
) -> Presentation:
    """Build an in-memory native presentation from a validated SceneSpec."""
    scene = validate_scene_spec(scene_spec)
    prs = Presentation(); prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
    canvas = scene["canvas"]
    for slide_spec in scene["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for node in slide_spec["nodes"]:
            if node["kind"] == "shape":
                _add_shape(slide, node, canvas)
            elif node["kind"] == "text":
                _add_text(slide, node, canvas)
            elif node["kind"] == "chart":
                _add_chart(slide, node, canvas)
            else:
                _add_image(slide, node, canvas, assets)
    return prs


def _candidate_gate(candidate: Path, deck: dict[str, Any]) -> dict[str, Any]:
    from .validator import validate_presentation
    return validate_presentation(candidate, deck)


def _path_contains_symlink(path: Path) -> bool:
    """Reject an existing symlink at the output leaf or any ancestor."""
    return any(candidate.is_symlink() for candidate in (path, *path.parents))


def render_deck(deck_spec: dict[str, Any], output: str | Path, *, protected_paths: Iterable[str | Path] = (), variants: dict[str, str] | None = None, slide_ids: Iterable[str] | None = None) -> dict[str, Any]:
    """Compile, render privately, validate, then atomically publish a native PPTX."""
    deck = validate_deck_spec(deck_spec)
    raw_output = Path(output).expanduser().absolute()
    if _path_contains_symlink(raw_output):
        raise RenderError("output path must not contain a symlink")
    output_path = raw_output
    resolved_output = output_path.resolve(strict=False)
    protected = {Path(path).resolve() for path in protected_paths}
    if resolved_output in protected:
        raise RenderError("output must not overwrite an input or protected path")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    if _path_contains_symlink(output_path):
        raise RenderError("output path must not contain a symlink")
    assets = _asset_map(deck)
    asset_paths = {Path(asset["path"]).resolve() for asset in deck["assets"]}
    asset_paths.update(Path(asset["fallback_path"]).resolve() for asset in deck["assets"] if asset["kind"] == "svg")
    if resolved_output in asset_paths:
        raise RenderError("output must not overwrite an asset path")
    scene = compile_deck(deck, variants=variants, slide_ids=slide_ids)
    selected_ids = {slide["slide_id"] for slide in scene["slides"]}
    selected_deck = copy.deepcopy(deck); selected_deck["slides"] = [slide for slide in deck["slides"] if slide["slide_id"] in selected_ids]
    prs = render_scene_presentation(scene, assets=assets)
    prs.core_properties.title = deck["title"]
    prs.core_properties.subject = "Generated by bounded AI PPTX Composer"
    fd, candidate_name = tempfile.mkstemp(prefix=f".{output_path.name}.", suffix=".candidate", dir=output_path.parent); os.close(fd)
    candidate = Path(candidate_name)
    try:
        prs.save(candidate)
        reopened = Presentation(candidate)
        if len(reopened.slides) != len(scene["slides"]):
            raise RenderError("candidate slide count mismatch")
        gate = _candidate_gate(candidate, selected_deck)
        if gate["status"] == "invalid":
            raise RenderError("candidate validation failed")
        if _path_contains_symlink(output_path):
            raise RenderError("output path must not contain a symlink")
        os.replace(candidate, output_path)
    except Exception:
        candidate.unlink(missing_ok=True)
        raise
    return {"status": "rendered", "deck_id": deck["deck_id"], "slide_count": len(scene["slides"]), "output": str(output_path), "sha256": hashlib.sha256(output_path.read_bytes()).hexdigest()}
