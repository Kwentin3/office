import tempfile
import unittest
import zipfile
from pathlib import Path

from pptx import Presentation

from pptx_ai_composer.renderer import render_deck
from pptx_ai_composer.validator import validate_presentation
from tests.fixtures import expanded_deck_spec, valid_deck_spec


class PresentationValidatorTests(unittest.TestCase):
    def setUp(self):
        self.tempdir = tempfile.TemporaryDirectory()
        self.root = Path(self.tempdir.name)
        self.deck = self.root / "deck.pptx"
        self.spec = valid_deck_spec()
        render_deck(self.spec, self.deck)

    def tearDown(self):
        self.tempdir.cleanup()

    def test_validates_structural_semantic_and_geometry_subset(self):
        report = validate_presentation(self.deck, self.spec)
        self.assertEqual(report["status"], "valid_with_unexecuted_gates")
        self.assertEqual(report["structural"]["status"], "pass")
        self.assertEqual(report["semantic"]["status"], "pass")
        self.assertEqual(report["geometry"]["status"], "pass")
        self.assertEqual(report["application"]["status"], "not_executed")
        self.assertEqual(report["visual"]["status"], "not_executed")
        self.assertIn("renderer unavailable", report["visual"]["reason"])

    def test_validates_expanded_archetypes_semantically(self):
        spec = expanded_deck_spec()
        output = self.root / "expanded.pptx"
        render_deck(spec, output)
        report = validate_presentation(output, spec)
        self.assertEqual(report["status"], "valid_with_unexecuted_gates")
        self.assertEqual(report["semantic"]["issues"], [])

    def test_detects_missing_source_citation(self):
        spec = valid_deck_spec()
        spec["slides"][0]["source_ids"] = ["src_1"]
        output = self.root / "with-cover-citation.pptx"
        render_deck(spec, output)
        prs = Presentation(output)
        for slide in prs.slides:
            for shape in list(slide.shapes):
                if shape.name == "scene:source.footer":
                    slide.shapes._spTree.remove(shape._element)
        prs.save(output)
        report = validate_presentation(output, spec)
        self.assertEqual(report["status"], "invalid")
        issues = report["semantic"]["issues"]
        self.assertTrue(any(issue["code"] == "missing_source_citation" for issue in issues))

    def test_detects_missing_expected_text(self):
        changed = self.spec.copy()
        changed["slides"] = [dict(slide) for slide in self.spec["slides"]]
        changed["slides"][0]["title"] = "A title that is not in the file"
        report = validate_presentation(self.deck, changed)
        self.assertEqual(report["status"], "invalid")
        self.assertTrue(any(issue["code"] == "missing_expected_text" for issue in report["semantic"]["issues"]))

    def test_detects_off_slide_shape(self):
        prs = Presentation(self.deck)
        prs.slides[0].shapes.add_textbox(prs.slide_width, 0, 1000, 1000).text = "outside"
        prs.save(self.deck)
        report = validate_presentation(self.deck, self.spec)
        self.assertEqual(report["status"], "invalid")
        self.assertTrue(any(issue["code"] == "off_slide_shape" for issue in report["geometry"]["issues"]))

    def test_rejects_duplicate_zip_member(self):
        import warnings
        with warnings.catch_warnings():
            warnings.simplefilter("ignore", UserWarning)
            with zipfile.ZipFile(self.deck, "a") as package:
                package.writestr("ppt/presentation.xml", b"duplicate")
        report = validate_presentation(self.deck, self.spec)
        self.assertEqual(report["status"], "invalid")
        self.assertTrue(any(issue["code"] == "duplicate_zip_member" for issue in report["structural"]["issues"]))


if __name__ == "__main__":
    unittest.main()
