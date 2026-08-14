import hashlib
import tempfile
import unittest
import zipfile
from pathlib import Path
from unittest.mock import patch
from PIL import Image
from pptx import Presentation

from pptx_ai_composer.compiler import compile_deck
from pptx_ai_composer.renderer import RenderError, render_deck, render_scene_presentation

from tests.fixtures import valid_deck_spec


class NativeRendererTests(unittest.TestCase):
    def setUp(self):
        self.tempdir = tempfile.TemporaryDirectory()
        self.root = Path(self.tempdir.name)
        self.output = self.root / "deck.pptx"

    def tearDown(self):
        self.tempdir.cleanup()

    def spec(self):
        return valid_deck_spec()

    def test_renders_native_text_shapes_and_editable_chart(self):
        result = render_deck(self.spec(), self.output)

        self.assertEqual(result["status"], "rendered")
        self.assertEqual(result["slide_count"], 3)
        prs = Presentation(self.output)
        self.assertEqual(len(prs.slides), 3)
        cover_text = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))
        self.assertIn("Hermes pilot", cover_text)
        self.assertTrue(any(shape.has_chart for shape in prs.slides[2].shapes))
        comparison_text = "\n".join(shape.text for shape in prs.slides[1].shapes if hasattr(shape, "text"))
        self.assertIn("Five tools", comparison_text)
        self.assertIn("Auditable actions", comparison_text)
        all_text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        self.assertNotIn("VERTICAL SLICE", all_text)
        self.assertNotIn("Sources: src_1", all_text)
        self.assertIn("Источник: Pilot workbook", all_text)
        with zipfile.ZipFile(self.output) as package:
            self.assertTrue(any(name.startswith("ppt/charts/chart") for name in package.namelist()))
            self.assertTrue(any(name.startswith("ppt/embeddings/") and name.endswith(".xlsx") for name in package.namelist()))

    def test_renders_png_asset_on_cover_and_binds_hash(self):
        image = self.root / "hero.png"
        Image.new("RGB", (640, 360), "#D95D39").save(image)
        digest = hashlib.sha256(image.read_bytes()).hexdigest()
        spec = self.spec()
        spec["assets"] = [{
            "asset_id": "hero",
            "kind": "png",
            "path": str(image),
            "sha256": digest,
            "alt_text": "Warm abstract field",
        }]
        spec["slides"][0]["asset_id"] = "hero"

        render_deck(spec, self.output)

        prs = Presentation(self.output)
        self.assertTrue(any(shape.shape_type == 13 for shape in prs.slides[0].shapes))

    def test_renders_svg_through_hash_bound_png_fallback(self):
        svg = self.root / "hero.svg"
        fallback = self.root / "hero.png"
        svg.write_text('<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100"><rect width="100" height="100" fill="#2A7F72"/></svg>', encoding="utf-8")
        Image.new("RGB", (400, 400), "#2A7F72").save(fallback)
        spec = self.spec()
        spec["assets"] = [{
            "asset_id": "hero",
            "kind": "svg",
            "path": str(svg),
            "sha256": hashlib.sha256(svg.read_bytes()).hexdigest(),
            "fallback_path": str(fallback),
            "fallback_sha256": hashlib.sha256(fallback.read_bytes()).hexdigest(),
            "alt_text": "Teal square",
        }]
        spec["slides"][0]["asset_id"] = "hero"
        render_deck(spec, self.output)
        prs = Presentation(self.output)
        self.assertTrue(any(shape.shape_type == 13 for shape in prs.slides[0].shapes))

    def test_cover_asset_variant_matrix_renders_png_and_svg_fallback(self):
        from pptx_ai_composer.library import get_catalog
        variants = get_catalog()["archetypes"]["cover"]["variants"]
        png = self.root / "matrix.png"
        svg = self.root / "matrix.svg"
        fallback = self.root / "matrix-fallback.png"
        Image.new("RGB", (120, 80), "#D95D39").save(png)
        svg.write_text('<svg xmlns="http://www.w3.org/2000/svg"><rect width="120" height="80" fill="#2A7F72"/></svg>', encoding="utf-8")
        Image.new("RGB", (120, 80), "#2A7F72").save(fallback)
        assets = [
            {"asset_id": "png", "kind": "png", "path": str(png), "sha256": hashlib.sha256(png.read_bytes()).hexdigest(), "alt_text": "PNG"},
            {"asset_id": "svg", "kind": "svg", "path": str(svg), "sha256": hashlib.sha256(svg.read_bytes()).hexdigest(), "fallback_path": str(fallback), "fallback_sha256": hashlib.sha256(fallback.read_bytes()).hexdigest(), "alt_text": "SVG"},
        ]
        for asset in assets:
            for variant in variants:
                with self.subTest(kind=asset["kind"], variant=variant):
                    deck = self.spec(); deck["assets"] = [asset]; deck["slides"][0]["asset_id"] = asset["asset_id"]
                    output = self.root / f"{asset['kind']}-{variant}.pptx"
                    result = render_deck(deck, output, variants={"s_cover": variant}, slide_ids=["s_cover"])
                    self.assertEqual(result["slide_count"], 1)
                    prs = Presentation(output)
                    self.assertTrue(any(shape.shape_type == 13 for shape in prs.slides[0].shapes))

    def test_refuses_svg_with_active_content(self):
        svg = self.root / "unsafe.svg"
        fallback = self.root / "fallback.png"
        svg.write_text('<svg xmlns="http://www.w3.org/2000/svg"><script>alert(1)</script></svg>', encoding="utf-8")
        Image.new("RGB", (100, 100), "white").save(fallback)
        spec = self.spec()
        spec["assets"] = [{
            "asset_id": "unsafe",
            "kind": "svg",
            "path": str(svg),
            "sha256": hashlib.sha256(svg.read_bytes()).hexdigest(),
            "fallback_path": str(fallback),
            "fallback_sha256": hashlib.sha256(fallback.read_bytes()).hexdigest(),
            "alt_text": "Unsafe SVG",
        }]
        spec["slides"][0]["asset_id"] = "unsafe"
        with self.assertRaisesRegex(RenderError, "unsafe svg"):
            render_deck(spec, self.output)

    def test_refuses_asset_when_bytes_do_not_match_contract_hash(self):
        image = self.root / "hero.png"
        Image.new("RGB", (100, 100), "red").save(image)
        spec = self.spec()
        spec["assets"] = [{
            "asset_id": "hero",
            "kind": "png",
            "path": str(image),
            "sha256": "0" * 64,
            "alt_text": "Red square",
        }]
        spec["slides"][0]["asset_id"] = "hero"
        with self.assertRaisesRegex(RenderError, "asset hash mismatch"):
            render_deck(spec, self.output)
        self.assertFalse(self.output.exists())

    def test_does_not_publish_when_candidate_validator_rejects(self):
        with patch("pptx_ai_composer.renderer._candidate_gate", return_value={"status": "invalid"}):
            with self.assertRaisesRegex(RenderError, "candidate validation failed"):
                render_deck(self.spec(), self.output)
        self.assertFalse(self.output.exists())

    def test_refuses_source_equals_output(self):
        spec_path = self.root / "deck.json"
        spec_path.write_text("{}", encoding="utf-8")
        with self.assertRaisesRegex(RenderError, "output must not overwrite"):
            render_deck(valid_deck_spec(), spec_path, protected_paths=[spec_path])

    def test_refuses_output_collision_with_asset_path(self):
        asset = self.root / "asset.png"
        Image.new("RGB", (40, 40), "red").save(asset)
        deck = valid_deck_spec()
        deck["assets"] = [{
            "asset_id": "hero", "kind": "png", "path": str(asset),
            "sha256": hashlib.sha256(asset.read_bytes()).hexdigest(), "alt_text": "Hero",
        }]
        deck["slides"][0]["asset_id"] = "hero"
        before = asset.read_bytes()
        with self.assertRaisesRegex(RenderError, "asset path"):
            render_deck(deck, asset)
        self.assertEqual(asset.read_bytes(), before)

    def test_scene_backend_builds_presentation_without_semantic_deck(self):
        scene = compile_deck(self.spec(), slide_ids=["s_compare"])
        prs = render_scene_presentation(scene, assets={})
        self.assertEqual(len(prs.slides), 1)
        text = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))
        self.assertIn("Five tools", text)
        self.assertIn("Pilot workbook", text)
        names = {shape.name for shape in prs.slides[0].shapes}
        self.assertIn("scene:comparison.left.card", names)

    def test_render_deck_accepts_named_variant_and_slide_selection(self):
        output = self.root / "single.pptx"
        result = render_deck(self.spec(), output, variants={"s_compare": "compact"}, slide_ids=["s_compare"])
        self.assertEqual(result["slide_count"], 1)
        prs = Presentation(output)
        self.assertEqual(len(prs.slides), 1)
        self.assertEqual(prs.slides[0].shapes[0].name, "scene:background")


if __name__ == "__main__":
    unittest.main()
