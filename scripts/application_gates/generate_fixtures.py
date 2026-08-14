"""Generate deterministic, non-sensitive fixtures for manual application gates."""

from __future__ import annotations

import argparse
import json
import os
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches


def generate(output_dir: Path) -> None:
    if not output_dir.is_absolute():
        raise ValueError("output directory must be absolute")
    output_dir.mkdir(mode=0o700)

    docx_path = output_dir / "basic.docx"
    document = Document()
    document.add_heading("Compatibility fixture", level=1)
    document.add_paragraph("Preserve this paragraph and formatting.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Key"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "answer"
    table.cell(1, 1).text = "42"
    document.save(docx_path)

    xlsx_path = output_dir / "recalculation.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Calc"
    sheet["A1"] = 1
    sheet["B1"] = "=A1+1"
    sheet["A2"] = "Preserve"
    workbook.save(xlsx_path)

    pptx_path = output_dir / "basic.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    shape.name = "slot:title"
    shape.text = "Compatibility fixture"
    presentation.save(pptx_path)

    manifest = {
        "schema_version": 1,
        "fixtures": [
            {"id": "docx-basic", "artifact_type": "docx", "path": docx_path.name},
            {"id": "xlsx-recalculation", "artifact_type": "xlsx", "path": xlsx_path.name},
            {"id": "pptx-basic", "artifact_type": "pptx", "path": pptx_path.name},
        ],
    }
    cases = {
        "schema_version": 1,
        "cells": [
            {
                "sheet": "Calc",
                "cell": "B1",
                "mode": "verify",
                "expected_type": "number",
                "expected_value": 2,
                "tolerance": 0,
            }
        ],
    }
    for path, payload in (
        (output_dir / "manifest.json", manifest),
        (output_dir / "recalculation-cases.json", cases),
    ):
        descriptor = os.open(path, os.O_WRONLY | os.O_CREAT | os.O_EXCL, 0o600)
        with os.fdopen(descriptor, "w", encoding="utf-8") as stream:
            json.dump(payload, stream, ensure_ascii=False, separators=(",", ":"))
            stream.write("\n")


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output-dir", required=True)
    args = parser.parse_args()
    try:
        generate(Path(args.output_dir))
    except (OSError, ValueError):
        return 2
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
